# This Python file uses the following encoding: iso-8859-15
#
# slides20.py -- 2020 refresh of slide template
#
import pprint
import inspect
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.util import Inches
from lxml import etree
import io
import os
import PIL
from PIL import ImageOps
from PIL import ImageDraw, ImageFont
import argparse
import json
import glob
import pytesseract
from datetime import datetime
import locale

import statistics
from statistics import mode

from pptx.oxml import parse_from_template, parse_xml
from pptx.oxml.dml.fill import CT_GradientFillProperties
from pptx.oxml.ns import nsdecls

import sqlite3
import dateutil.parser as parser


# If you don't have tesseract executable in your PATH, include the following:
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract'

#
# Globals
#

assetRoot = "assets/"

#
# Helper functions
#

from shs2phss import shs2phss

def get_song_paths(book, number):
	# map shs songs to phss files
	if book == 'shs':
		if number == 34:
			number = 109
			book = "eh"
		elif str(number) in shs2phss:
			number = shs2phss[str(number)]
			book = "phss"
	song = f"{number:03d}"
	pathname = '/'.join(['ehsf', book, song])
	basename = pathname + "/" + book + "-" + song
	rawname = pathname + "/raw/" + book + "-" + song
	return song, pathname, basename, rawname


def get_song_paths_new(book, number):
	paths = dict()
	# map shs songs to phss files
	if book == 'shs':
		if number == 34:
			number = 109
			book = "eh"
		if str(number) in shs2phss:
			number = shs2phss[str(number)]
			book = "phss"
	song = f"{number:03d}"
	paths['engpath'] = '/'.join(['ehsf', book, song])
	paths['engbase'] = paths['engpath'] + "/" + book + "-" + song
	paths['raw'] = paths['engpath'] + "/raw/" + book + "-" + song
	paths['esppath'] = '/'.join(['ehsf', 'esp', book, song])
	paths['espbase'] = paths['esppath'] + "/" + book + "-" + song
	return song, paths


def fitText(frame, font_family, max_size, bold, step_size=12, file=None, features=None, spacing=0):
	while True:
		try:
			size = frame.fit_text(font_family, max_size, bold, font_file=file, features=features, spacing=spacing)
			return size
		except Exception as e:
			max_size = max_size - step_size
			print("ex", max_size, font_family, e)
			if max_size < 12:
				return max_size


def set_placeholder_size(v, top, left, width, height):
	v.top = Inches(top)
	v.left = Inches(left)
	v.width = Inches(width)
	v.height = Inches(height)

# quirky way we have to set position of an image in a placeholder
# must set all four values for both placeholder and pic
def set_placeholder_pic(v, img, top, left, width, height):
	set_placeholder_size(v, top, left, width, height)
	pic = v.insert_picture(img)
	pic.top = Inches(top)
	pic.left = Inches(left)
	pic.width = Inches(width)
	pic.height = Inches(height)

def get_placeholder(slide, pid):
	for v in slide.placeholders:
#		print(v.name, v.shape_id, v.top.inches, v.left.inches)
		if v.shape_id == pid:
			return v
	return None

def dump_placeholders(slide):
	for v in slide.placeholders:
		print(v.name, v.shape_id, v.top.inches, v.left.inches)


##############################################################################
##
## PFTL Song Image Processing
##
##############################################################################

def analyze_image(picture):
	img = PIL.Image.open(picture)
	print(img.mode, img.size)
	if img.mode != "1":
		img = img.convert(mode="1")
	print(img.mode, img.size)
	pixels = list(img.getdata())
	width, height = img.size
	top = -1
	bot = height + 1
	left = width + 1
	right = -1
	staff = -1
	for ny in range(height):
		row = pixels[(ny * width):(ny * width)+width]
		np = row.count(255)
		if np != width:
			if top == -1:
				top = ny
			bot = ny
			if np < width / 2:
				if staff == -1:
					staff = ny
			tl = row.index(0)
			if tl < left:
				left = tl
			rt = width - row[::-1].index(0)
			if rt > right:
				right = rt
	print(top, bot, staff, left, right)
	return dict(width=width, height=height, top=top/height, bot=bot/height, staff=staff/height, left=left/width, right=right/width)



# Determine limiting axis (based on > or < 16:9)
# Extract the source window
# Add % padding (based on pixel dimension of source window)
# Determine pixel size of other axis
# Create blank canvas and center cropped portion

# set_window - determine PowerPoint window size and position
#  window = dimensions in % of source window
#  padding = padding factor (ex, 1.05) to add
#  iar = aspect ratio of input images
#  returns [window dimensions in inches for PowerPoint]

# Original: 9.6 x 5.4 (16x9)
# 
# < 1.25 (max height, narrow width)
# Tall: 7.5 x 6 (5x4) 1.250
# 
#   < 1.388 (max width, lower height)
# Mid: 7.5 x 5.4  (5x3.6) 1.388  
#   > 1.388 (max height, narrow width)
# 
# Wide: 9 x 5.4 (5x3) 1.677
# > 1.67 (max width, lower height)

def set_window(window, padding, iar):
	print(window)
	wd = window[2] * padding
	ht = window[3] * padding
	ar = (wd/ht) * iar
	print(wd, ht)
	print(ar)
	if ar <= (7.5/5.4):
		orientation = "tall"
		if ar <= (5.0/4.0):
			# tall case: scale to max h, adjust width
			pp_top = 0.125
			pp_height = 6.0
			pp_width = 6.0 * ar
			pp_left = 0.80 + ((7.5 - pp_width) / 2)
		else:
			# wide case: scale to maxw, adjust height
			pp_left = 0.80
			pp_width = 7.5
			pp_height = 7.5 / ar
			pp_top = 0.125 + ((6.0 - pp_height) / 2)
	else:
		orientation = "wide"
		if ar <= (5.0/3.0):
			# tall case: scale to max h, adjust width
			pp_top = 0.125
			pp_height = 5.4
			pp_width = 5.4 * ar
			pp_left = 0.80 + ((9.0 - pp_width) / 2)
		else:
			# wide case: scale to maxw, adjust height
			pp_left = 0.80
			pp_width = 9.0
			pp_height = 9.0 / ar
			pp_top = 0.125 + ((5.4 - pp_height) / 2)
	return orientation, [pp_top, pp_left, pp_width, pp_height]


# analyze all the images to determine crop for each
def set_crop_window(crop, meta):
	mtop = 1.1
	mbot = -0.1
	ml = 1.1
	mr = -0.1
	for key, cr in crop.items():
		mtop = min(mtop, cr["top"])
		mbot = max(mbot, cr["bot"])
		ml = min(ml, cr["left"])
		mr = max(mr, cr["right"])
		iar = cr["width"] / cr["height"]

	window = [ mtop, ml, mr - ml, mbot - mtop]
	padding = 1.05

	meta['window_orientation'], meta['window'] = set_window(window, padding, iar)
	return window, padding



def size_image_to_window(picture, crop, window, padding, basename, rawname, ndx):
	filename = basename + "-" + f"{ndx:02d}" + ".png"
	img = PIL.Image.open(picture)
	if rawname is not None:
		img.save(rawname + "-" + f"{ndx:02d}" + ".png")	# save the raw image
	if img.mode != "RGB":
		img = img.convert(mode="RGB")
# window = [ mtop, ml, mr - ml, mbot - mtop]
	mw = int((window[2] * crop["width"]) + 0.5)
	mh = int((window[3] * crop["height"]) + 0.5)
	left = window[1] * crop["width"]
	upper = window[0] * crop["height"]
	right = left + mw
	lower = upper + mh
	box = (left, upper, right, lower)
	print(box)
	img = img.crop(box)

	pw = mw * padding
	ph = mh * padding
	dx = (pw - mw) / 2
	dy = (ph - mh) / 2

	out = PIL.Image.new("RGB", (int(pw + 0.5), int(ph + 0.5)), color="white")
	out.paste(img, box=(int(dx), int(dy)))
	out.save(filename)



##############################################################################
##
## "Praise For The Lord" Parser (works on PPTX files)
##
##############################################################################

# Need to adjust for chorus!
#	Text: Lyrics Verse (for first slide of verse)
#	Picture
#	Text: 'Verse/chorus - Title'
#	Text: Credits (for first slide of verse)
#	Text: Copyright
#	Text: Song Number

def is_pftl_copyright(text):
	if text.find("Paperless") == -1:
		return False
	return True

def is_pftl_number(text):
	if len(text) <= 3 and text.isnumeric():
		return True
	return False

def is_pftl_title(text):
	print("title", text[0:1].isnumeric(), text[0:1] == 'c', len(text.splitlines()))
	if text[0:1].isnumeric() or text[0:1] == 'c' or text == 'Amen' or text[0:1].lower() == 's':
		if len(text.splitlines()) == 1:
			return True
	return False

def is_pftl_lyric(text):
	if text[0:1].isnumeric() and len(text.splitlines()) > 1:
		return True
	return False


def process_pftl_song(number):
	song, pathname, basename, rawname = get_song_paths("pftl", number)
	prs = Presentation("ehsf/pftl/pptx/" + song + '.pptx')
#	pprint.pprint(inspect.getmembers(prs))

	# Create a directory to hold our output
	os.makedirs(pathname, exist_ok=True)
	os.makedirs(pathname + "/raw", exist_ok=True)

	slides = prs.slides
#	pprint.pprint(inspect.getmembers(slides))

	crop = dict()

	meta = dict()
	meta['lyrics'] = dict()
	meta['verses'] = dict()
	meta['chorus'] = dict()
	meta['codas'] = dict()

	nVerse = 0

	for ndx, slide in enumerate(slides, 1):
		lyrics = None
		credit_lyric = None
		chorus = False
		sanctus = False
#		pprint.pprint(inspect.getmembers(slide))
#		print(etree.tostring(slide.element, pretty_print=True))
		print("\n")
		for ndy, shape in enumerate(slide.shapes, 1):
#			pprint.pprint(dir(shape))
			print(ndy, shape.shape_type)
			if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
				print(shape.placeholder_format.type)
				if shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
					crop[ndx] = analyze_image(io.BytesIO(shape.image.blob))
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				crop[ndx] = analyze_image(io.BytesIO(shape.image.blob))
#				pprint.pprint(inspect.getmembers(shape.image))
			if shape.has_text_frame:
				text_frame = shape.text_frame
				# some songs have blank text_frames embedded
				if len(text_frame.text) == 0:
					continue
				if is_pftl_number(text_frame.text):
					if 'number' not in meta:
						meta['number'] = text_frame.text
				elif is_pftl_copyright(text_frame.text):
					if 'copyright' not in meta:
						meta['copyright'] = text_frame.text
				elif is_pftl_title(text_frame.text):
					if 'title' not in meta:
						meta['title'] = text_frame.text[4:]
					if text_frame.text[0:1].isnumeric():
						thisVerse = int(text_frame.text.split()[0])
						if thisVerse != nVerse:
							nVerse = thisVerse
							meta['verses'][nVerse] = []
							if lyrics:
								meta['lyrics'][nVerse] = [ lyrics ]
						meta['verses'][nVerse].append(ndx)
					elif text_frame.text[0:1] == 'c':
						if nVerse not in meta['chorus']:
							meta['chorus'][nVerse] = []
						meta['chorus'][nVerse].append(ndx)
						chorus = True
					elif text_frame.text == 'Amen':
						if 'Amen' not in meta['codas']:
							meta['codas']['Amen'] = []
						meta['codas']['Amen'].append(ndx)
					elif text_frame.text[0:1].lower() == 's':
						if 'Sanctus' not in meta['codas']:
							meta['codas']['Sanctus'] = []
						meta['codas']['Sanctus'].append(ndx)
						sanctus = True
				elif is_pftl_lyric(text_frame.text):
					lyrics = text_frame.text
				else:
					credit_lyric = text_frame.text
				print(nVerse, text_frame.text, shape.shape_type)

		if credit_lyric:
			if chorus:
				vdx = str(nVerse) + "c"
				if vdx not in meta['lyrics']:
					meta['lyrics'][vdx] = []
				meta['lyrics'][vdx].append(credit_lyric)
			elif sanctus:
				if 'sanctus' not in meta['lyrics']:
					meta['lyrics']['sanctus'] = []
				meta['lyrics']['sanctus'].append(credit_lyric)
			else:
				if 'credits' not in meta:
					meta['credits'] = credit_lyric

	# analyze all the images to determine crop for each
	window, padding = set_crop_window(crop, meta)

	for ndx, slide in enumerate(slides, 1):
		for ndy, shape in enumerate(slide.shapes, 1):
			if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE) or ((shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER) and (shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT)):
				size_image_to_window(io.BytesIO(shape.image.blob), crop[ndx], window, padding, basename, rawname, ndx)

	# Output JSON data
	pprint.pprint(meta)
	with open(basename + ".json", 'w') as jsonfile:
		json.dump(meta, jsonfile, ensure_ascii=False, indent=4)


##############################################################################
##
## "Psalms, Hymns, and Spiritual Songs" Parser (works on exported PPTX files)
##
##############################################################################

def phss_get_images(slides, ndx, token, shape=0):
	files = []
	while True:
		print(ndx)
		if ndx < len(slides):
			if slides[ndx].shapes[shape].has_text_frame:
				print(slides[ndx].shapes[shape].text_frame.text)
				if slides[ndx].shapes[shape].text_frame.text.find(token) != -1:
					files.append(ndx)
					ndx = ndx + 1
					continue
		break
	return files, ndx

def process_phss_song_ppt(number):
	song, pathname, basename, rawname = get_song_paths("phss", number)
	prs = Presentation("ehsf/phss/pptx/" + song + '.pptx')
	slides = prs.slides

#	for ndx, slide in enumerate(slides, 1):
#		print(ndx)
#		print("\n")
#		for ndy, shape in enumerate(slide.shapes, 1):
#			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
##				crop[ndx] = analyze_image(io.BytesIO(shape.image.blob))
#				print("image")
#			if shape.has_text_frame:
#				text_frame = shape.text_frame
#				pprint.pprint(text_frame.text)

	# Create a directory to hold our output
	os.makedirs(pathname, exist_ok=True)
	os.makedirs(pathname + "/raw", exist_ok=True)

	meta = dict()
	meta['lyrics'] = dict()
	meta['verses'] = dict()
	meta['chorus'] = dict()
	meta['codas'] = dict()

	meta['number'] = str(number)
	meta['copyright'] = u"\xa9 2014 Sumphonia Productions, LLC"

	nVerse = 0
	nChorus = 0
	nCoda = 0

	# Extract lyrics from XML file
	with open("ehsf/phss/phss.xml", 'r') as xml:
		tree = etree.parse(xml)
	hymn = tree.xpath('/Hymnal/HymnEntry[@HymnNumber="' + str(number) + '"]')

	chorus_token = "_cho_"

	for element in hymn[0].iter():
		print("%s - %s" % (element.tag, element.text), element.attrib)
		if element.tag == "Title":
			meta['title'] = element.text
		elif element.tag == "HymnElement":
			if "HymnElementCategory" in element.attrib:
				if element.attrib["HymnElementCategory"] == "Verse":
					nVerse = nVerse + 1
					meta['lyrics'][nVerse] = []
					for item in element.iter():
						if item.tag == "HymnLine":
							meta['lyrics'][nVerse].append(item.text)
				elif element.attrib["HymnElementCategory"] == "Chorus" or element.attrib["HymnElementCategory"] == "Refrain":
					nChorus = nChorus + 1
					meta['lyrics']['chorus'] = []
					for item in element.iter():
						if item.tag == "HymnLine":
							meta['lyrics']['chorus'].append(item.text)
					if element.attrib["HymnElementCategory"] == "Refrain":
						chorus_token = "_ref_"
				elif element.attrib["HymnElementCategory"] == "Coda":
					nCoda = nCoda + 1
					meta['lyrics']['Coda'] = []
					for item in element.iter():
						if item.tag == "HymnLine":
							meta['lyrics']['Coda'].append(item.text)


	# Extract credits from first slide in PPT file
	credits = []
	for ndy, shape in enumerate(prs.slides[0].shapes, 1):
		if shape.has_text_frame:
			text = shape.text_frame.text
			if text != meta['title'] and text != meta['number']:
				credits.append(text)
	meta['credits'] = '\n'.join(credits)

	# Generate list of PNG files and assign to verses
	png = 0
	ndx = 1
	pngfiles = dict()
	for verse in range(1, nVerse + 1):
		meta['verses'][verse] = []
		files, ndx = phss_get_images(prs.slides, ndx, "_st" + str(verse) + "_")
		for file in files:
			png = png + 1
			pngfiles[png] = file
			meta['verses'][verse].append(png)
		if nChorus != 0:
			meta['chorus'][verse] = []
			files, ndx = phss_get_images(prs.slides, ndx, chorus_token)
			for file in files:
				png = png + 1
				pngfiles[png] = file
				meta['chorus'][verse].append(png)
	if nCoda != 0:
		meta['codas']['Coda'] = []
		files, ndx = phss_get_images(prs.slides, ndx, "_coda_")
		for file in files:
			png = png + 1
			pngfiles[png] = file
			meta['codas']['Coda'].append(png)

	# Analyze the extent of all the images
	crop = dict()
	for ndx, slidendx in pngfiles.items():
		slide = prs.slides[slidendx]
		for ndy, shape in enumerate(slide.shapes, 1):
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				crop[ndx] = analyze_image(io.BytesIO(shape.image.blob))
		
	window, padding = set_crop_window(crop, meta)

	for ndx, slidendx in pngfiles.items():
		slide = prs.slides[slidendx]
		for ndy, shape in enumerate(slide.shapes, 1):
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				size_image_to_window(io.BytesIO(shape.image.blob), crop[ndx], window, padding, basename, rawname, ndx)
		
	# Output JSON data
	pprint.pprint(meta)
	with open(basename + ".json", 'w') as jsonfile:
		json.dump(meta, jsonfile, ensure_ascii=False, indent=4)


# take a PHSS standalone file (number.pptx in eh/pptx) and make an eh song from it
def process_phss_to_eh(number):
	song, pathname, basename, rawname = get_song_paths("eh", number)
	prs = Presentation("ehsf/eh/pptx/" + song + '.pptx')
	slides = prs.slides

	for ndx, slide in enumerate(slides, 1):
		print(ndx)
		print("\n")
		for ndy, shape in enumerate(slide.shapes, 1):
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#				crop[ndx] = analyze_image(io.BytesIO(shape.image.blob))
				print("image")
			if shape.has_text_frame:
				text_frame = shape.text_frame
				pprint.pprint(text_frame.text)

	# Create a directory to hold our output
	os.makedirs(pathname, exist_ok=True)
	os.makedirs(pathname + "/raw", exist_ok=True)

	meta = dict()
	meta['verses'] = dict()
	meta['chorus'] = dict()
	meta['codas'] = dict()

	meta['number'] = str(number)
	meta['copyright'] = u"\xa9 2014 Sumphonia Productions, LLC"

	# hardwired for now
	chorus_token = "Refrain"
	nVerse = 4
	nChorus = 0
	nCoda = 0

	# Extract credits from first slide in PPT file
	credits = []
	for ndy, shape in enumerate(prs.slides[0].shapes, 1):
		if shape.has_text_frame:
			text = shape.text_frame.text
			credits.append(text)
	meta['credits'] = '\n'.join(credits)

	print(meta['credits'])
	# how to get title? is it the first one?

	# Generate list of PNG files and assign to verses
	png = 0
	ndx = 1
	pngfiles = dict()
	for verse in range(1, nVerse + 1):
		meta['verses'][verse] = []
		files, ndx = phss_get_images(prs.slides, ndx, "Vs " + str(verse), 1)
		for file in files:
			png = png + 1
			pngfiles[png] = file
			meta['verses'][verse].append(png)
		if nChorus != 0:
			meta['chorus'][verse] = []
			files, ndx = phss_get_images(prs.slides, ndx, chorus_token, 1)
			for file in files:
				png = png + 1
				pngfiles[png] = file
				meta['chorus'][verse].append(png)
	if nCoda != 0:
		meta['codas']['Coda'] = []
		files, ndx = phss_get_images(prs.slides, ndx, "_coda_", 1)
		for file in files:
			png = png + 1
			pngfiles[png] = file
			meta['codas']['Coda'].append(png)

	# Analyze the extent of all the images
	crop = dict()
	for ndx, slidendx in pngfiles.items():
		slide = prs.slides[slidendx]
		for ndy, shape in enumerate(slide.shapes, 1):
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				crop[ndx] = analyze_image(io.BytesIO(shape.image.blob))
		
	window, padding = set_crop_window(crop, meta)

	for ndx, slidendx in pngfiles.items():
		slide = prs.slides[slidendx]
		for ndy, shape in enumerate(slide.shapes, 1):
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				size_image_to_window(io.BytesIO(shape.image.blob), crop[ndx], window, padding, basename, rawname, ndx)
		
	# Output JSON data
	pprint.pprint(meta)
	with open(basename + ".json", 'w') as jsonfile:
		json.dump(meta, jsonfile, ensure_ascii=False, indent=4)






##############################################################################
##
## Song Slide Generation
##
##############################################################################

#
# Generate output PowerPoint functions
#

def get_navlist(wp, png, lastpng):
	nd = []
	check_skip = False
	for p in wp:
		color = "gray"	# default color
		if 'pngs' in p:
			if png in p['pngs']:
				color = "blue"
				if 'skip' in p:
					if p['skip'] == 1:
						if png == lastpng:
							check_skip = True
							color = "highlight"	# drw: next is an arrow
			elif check_skip:
				color = "red"
				check_skip = False
		nd.append({ "text": p['text'], "color": color})
	return nd


def get_navigation(meta, verses, chorus, coda):
	# Generate master list of waypoints
	wp = []
	skipping = True
	for verse, pngs in meta['verses'].items():
		if verses and int(verse) not in verses:
			if not skipping:
				wp[-1]["skip"] = 1
				wp.append({ "text": "arrow"})
				skipping = True
			continue
		wp.append({ "text": verse, "pngs": pngs})
		skipping = False
		if verse in meta['chorus'].keys():
			if chorus and int(verse) not in chorus:
				if not skipping:
					wp[-1]["skip"] = 1
					wp.append({ "text": "arrow"})
					skipping = True
				continue
			wp.append({ "text": "C", "pngs": meta['chorus'][verse]})
			skipping = False
	if coda and len(meta['codas']) > 0:
		for verse, pngs in meta['codas'].items():
			wp.append({ "text": verse, "pngs": pngs})
		skipping = False

	if skipping:
		wp.pop()
		wp[-1]["skip"] = 0
#	pprint.pprint(wp)

	# Generate colorized list of waypoints for each png
	nav = dict()
	for verse, pngs in meta['verses'].items():
		if verses and int(verse) not in verses:
			continue
		for png in pngs:
			nav[png] = get_navlist(wp, png, pngs[-1])
		if verse in meta['chorus'].keys():
			if chorus and int(verse) not in chorus:
				continue
			for png in meta['chorus'][verse]:
				nav[png] = get_navlist(wp, png, meta['chorus'][verse][-1])
	if coda and len(meta['codas']) > 0:
		for verse, pngs in meta['codas'].items():
			for png in pngs:
				nav[png] = get_navlist(wp, png, pngs[-1])


	# Generate navinfo dict to return
	navinfo = dict()
	navinfo['nav'] = nav

	# 4.75 to 8 inches
	# large (0), can fit six
	# small (1), can fit ten
	numslots = len(wp)
	navinfo['size'] = 0 if numslots <= 6 else 1

#	pprint.pprint(navinfo)
	return navinfo


# at least one of verses/chorus must be set

vex_eng = dict(omit="Omit Verse ", omits="Omit Verses ", chorusafterone="Chorus After Verse ", chorusaftermany="Chorus After Verses ", chorusafterlast="Chorus After Last Verse", nochorus="Omit the Chorus")
vex_esp = dict(omit="Omitir el verso ", omits="Omitir los versos ", chorusafterone="Coro después del verso ", chorusaftermany="Coro después de los versos ", chorusafterlast="Coro después del último verso", nochorus="Omitir el coro")

def get_verse_exceptions(meta, verses, chorus, repeat, vex):
	if repeat:
		return "Repeat Song"
	elif verses is None:
		# we are doing all verses
		sing = []
		for verse in meta['chorus'].keys():
			if int(verse) in chorus:
				sing.append(verse)
		if len(sing) == 1:
			if int(sing[0]) == len(meta['verses']):
				return vex['chorusafterlast']
			else:
				return vex['chorusafterone'] + sing[0]
		elif len(sing) == 2:
			return vex['chorusaftermany'] + sing[0] + " & " + sing[1]
		elif len(sing) == 0:
			return vex['nochorus']
		else:
			return vex['chorusaftermany'] + ', '.join(sing)
	else:
		skipped = len(meta['verses']) - len(verses)
		if skipped < len(verses):
			omit = []
			for verse in meta['verses'].keys():
				if int(verse) not in verses:
					omit.append(verse)
			if len(omit) == 1:
				return vex['omit'] + omit[0]
			elif len(omit) == 2:
				return vex['omits'] + omit[0] + " & " + omit[1]
			else:
				return vex['omits'] + ', '.join(omit)
		else:
			sing = []
			for verse in meta['verses'].keys():
				if int(verse) in verses:
					sing.append(verse)
			if len(sing) == 1:
				return "Verse " + sing[0] + " Only"
			elif len(sing) == 2:
				return "Verses " + sing[0] + " & " + sing[1]
			else:
				return "Verses " + ', '.join(sing)
	return "All Verses"


MASTER_TITLE = 0
MASTER_STATIC = 1

# English only slide
LAYOUT_TITLE_ENG = 4
LAYOUT_TITLE_ENG_BACKGROUND = 2
LAYOUT_TITLE_ENG_NAVBAR_BG = 3
LAYOUT_TITLE_ENG_CIRCLE_BG = 4
LAYOUT_TITLE_ENG_TITLE = 5
LAYOUT_TITLE_ENG_DETAIL = 6
LAYOUT_TITLE_ENG_QUOTE = 7
LAYOUT_TITLE_ENG_CREDITS = 8
LAYOUT_TITLE_ENG_REFERENCE = 9
LAYOUT_TITLE_ENG_LOGO = 10
LAYOUT_TITLE_ENG_HIGHLIGHT = 11
LAYOUT_TITLE_ENG_NAVBAR = 12
LAYOUT_TITLE_ENG_CALLOUT = 13
LAYOUT_TITLE_ENG_COPYRIGHT = 14
LAYOUT_TITLE_ENG_LEADER = 15

# Bilingual
LAYOUT_TITLE_BIL = 5
LAYOUT_TITLE_BIL_SLOW = 6
LAYOUT_TITLE_BIL_BACKGROUND = 2
LAYOUT_TITLE_BIL_NAVBAR_BG = 3
LAYOUT_TITLE_BIL_CIRCLE_BG = 4
LAYOUT_TITLE_BIL_ENG_TITLE = 5
LAYOUT_TITLE_BIL_ENG_DETAIL = 6
LAYOUT_TITLE_BIL_ENG_QUOTE = 7
LAYOUT_TITLE_BIL_ENG_CREDITS = 8
LAYOUT_TITLE_BIL_ENG_REFERENCE = 9
LAYOUT_TITLE_BIL_ENG_LOGO = 10
LAYOUT_TITLE_BIL_ENG_HIGHLIGHT = 11
LAYOUT_TITLE_BIL_ENG_NAVBAR = 12
LAYOUT_TITLE_BIL_ENG_CALLOUT = 13
LAYOUT_TITLE_BIL_ENG_COPYRIGHT = 14
LAYOUT_TITLE_BIL_LEADER = 15
LAYOUT_TITLE_BIL_ESP_TITLE = 16
LAYOUT_TITLE_BIL_ESP_DETAIL = 17
LAYOUT_TITLE_BIL_ESP_QUOTE = 18
LAYOUT_TITLE_BIL_ESP_CREDITS = 19
LAYOUT_TITLE_BIL_ESP_REFERENCE = 20
LAYOUT_TITLE_BIL_ESP_LOGO = 21
LAYOUT_TITLE_BIL_ESP_HIGHLIGHT = 22
LAYOUT_TITLE_BIL_ESP_NAVBAR = 23
LAYOUT_TITLE_BIL_ESP_CALLOUT = 24
LAYOUT_TITLE_BIL_ESP_COPYRIGHT = 25


# Generic title slide indices
LAYOUT_TITLE_BACKGROUND = 0
LAYOUT_TITLE_NAVBAR_BG = 1
LAYOUT_TITLE_CIRCLE_BG = 2
LAYOUT_TITLE_TITLE = 3
LAYOUT_TITLE_DETAIL = 4
LAYOUT_TITLE_QUOTE = 5
LAYOUT_TITLE_CREDITS = 6
LAYOUT_TITLE_REFERENCE = 7
LAYOUT_TITLE_LOGO = 8
LAYOUT_TITLE_HIGHLIGHT = 9
LAYOUT_TITLE_NAVBAR = 10
LAYOUT_TITLE_CALLOUT = 11
LAYOUT_TITLE_COPYRIGHT = 12
LAYOUT_TITLE_LEADER = 13

ndx2ben = [
	LAYOUT_TITLE_BIL_BACKGROUND,
	LAYOUT_TITLE_BIL_NAVBAR_BG,
	LAYOUT_TITLE_BIL_CIRCLE_BG,
	LAYOUT_TITLE_BIL_ENG_TITLE,
	LAYOUT_TITLE_BIL_ENG_DETAIL,
	LAYOUT_TITLE_BIL_ENG_QUOTE,
	LAYOUT_TITLE_BIL_ENG_CREDITS,
	LAYOUT_TITLE_BIL_ENG_REFERENCE,
	LAYOUT_TITLE_BIL_ENG_LOGO,
	LAYOUT_TITLE_BIL_ENG_HIGHLIGHT,
	LAYOUT_TITLE_BIL_ENG_NAVBAR,
	LAYOUT_TITLE_BIL_ENG_CALLOUT,
	LAYOUT_TITLE_BIL_ENG_COPYRIGHT,
	LAYOUT_TITLE_BIL_LEADER
]

ndx2bes = [
	LAYOUT_TITLE_BIL_BACKGROUND,
	LAYOUT_TITLE_BIL_NAVBAR_BG,
	LAYOUT_TITLE_BIL_CIRCLE_BG,
	LAYOUT_TITLE_BIL_ESP_TITLE,
	LAYOUT_TITLE_BIL_ESP_DETAIL,
	LAYOUT_TITLE_BIL_ESP_QUOTE,
	LAYOUT_TITLE_BIL_ESP_CREDITS,
	LAYOUT_TITLE_BIL_ESP_REFERENCE,
	LAYOUT_TITLE_BIL_ESP_LOGO,
	LAYOUT_TITLE_BIL_ESP_HIGHLIGHT,
	LAYOUT_TITLE_BIL_ESP_NAVBAR,
	LAYOUT_TITLE_BIL_ESP_CALLOUT,
	LAYOUT_TITLE_BIL_ESP_COPYRIGHT,
	LAYOUT_TITLE_BIL_LEADER
]


LAYOUT_WELCOME = 7
LAYOUT_WELCOME_BIL = 8

LAYOUT_WELCOME_BACKGROUND = 2
LAYOUT_WELCOME_NAVBAR_BG = 3
LAYOUT_WELCOME_ENG_LOGO = 4
LAYOUT_WELCOME_ENG_ORDER = 5
LAYOUT_WELCOME_ENG_HIGHLIGHT = 6
LAYOUT_WELCOME_ENG_NAVBAR = 7
LAYOUT_WELCOME_ENG_DATE = 8
LAYOUT_WELCOME_ENG_SILENCE = 9
LAYOUT_WELCOME_LEADERS = 10
LAYOUT_WELCOME_ESP_LOGO = 11
LAYOUT_WELCOME_ESP_ORDER = 12
LAYOUT_WELCOME_ESP_HIGHLIGHT = 13
LAYOUT_WELCOME_ESP_NAVBAR = 14
LAYOUT_WELCOME_ESP_DATE = 15
LAYOUT_WELCOME_ESP_SILENCE = 16

# Generic welcome slide indices
LAYOUT_WELCOME_NDX_LOGO = 0
LAYOUT_WELCOME_NDX_ORDER = 1
LAYOUT_WELCOME_NDX_HIGHLIGHT = 2
LAYOUT_WELCOME_NDX_NAVBAR = 3
LAYOUT_WELCOME_NDX_DATE = 4
LAYOUT_WELCOME_NDX_SILENCE = 5



phlist = { "eng": [LAYOUT_TITLE_TITLE, LAYOUT_TITLE_DETAIL, LAYOUT_TITLE_QUOTE, LAYOUT_TITLE_CREDITS, LAYOUT_TITLE_REFERENCE, LAYOUT_TITLE_CALLOUT, LAYOUT_TITLE_COPYRIGHT],
		   "esp": [LAYOUT_TITLE_TITLE, LAYOUT_TITLE_DETAIL, LAYOUT_TITLE_QUOTE, LAYOUT_TITLE_CREDITS, LAYOUT_TITLE_REFERENCE, LAYOUT_TITLE_CALLOUT, LAYOUT_TITLE_COPYRIGHT],
		   "ldr": [LAYOUT_TITLE_LEADER] }

bookColors = { "pftl": RGBColor(106, 0, 8), "shs": RGBColor(0, 0, 0), "phss": RGBColor(5, 29, 69), "eh": RGBColor(142, 172, 208) }

eng_fonts = [["Avenir Next LT Pro", "AvenirNextLTPro-Regular.ttf"], ["Avenir Next LT Pro", "AvenirNextLTPro-Bold.otf"]]
esp_fonts = [["Alegreya Sans Medium", "AlegreyaSans-Medium.otf"], ["Alegreya Sans ExtraBold", "AlegreyaSans-ExtraBold.ttf"]]


def get_background_fn(default, item, blur, zoom):
	bgname = default
	if item is not None:
		if 'background' in item:
			bgname = item['background']
		elif 'meta' in item:
			if 'background' in item['meta']:
				bgname = item['meta']['background']
	basefn = bgname + "/" + bgname
	if blur:
		basefn = basefn + "-blur"
	if zoom:
		basefn = basefn + "-zoom"
	return assetRoot + "backgrounds/" + basefn + ".jpg"

def get_background(default, item=None, blur=False, zoom=False):
	fn = get_background_fn(default, item, blur, zoom)
	if blur or zoom:
		if not os.path.exists(fn):
			fn = get_background_fn(default, item, False, False)
	return fn

def get_song_bubbles(meta, verses, chorus, coda):
	wp = []
	for verse, pngs in meta['verses'].items():
		if verses and int(verse) not in verses:
			wp.append({ "bubble": str(verse), "sing": False })
		else:
			wp.append({ "bubble": str(verse), "sing": True })

		if chorus:
			if verse in meta['chorus'].keys():
				if int(verse) in chorus:
					wp.append({ "bubble": "C", "sing": True })

	if coda and len(meta['codas']) > 0:
		for verse, pngs in meta['codas'].items():
			text = verse	# default
			if text == "Coda":
				text = "C"	# tbd: use coda icon
			if text == "Sanctus":
				text = "S"
			if text == "Amen":
				text = "A"
			wp.append({ "bubble": text, "sing": True})

	pprint.pprint(wp)
	return wp



def add_bubble(slide, top, xleft, slotndx, text, sing):
	diameter = 0.4
	slot_width = 0.44
#	slot_width = 0.48
	fontsize = 16
	line_spacing = 0.9

	left = xleft + (slot_width * slotndx)

	shapes = slide.shapes
	shape = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(diameter), Inches(diameter))

	#1B255E -- eh dark blue
	#ED5959 -- eh red higlight

	fill = shape.fill
	fill.solid()
	if sing:
		fill.fore_color.rgb = RGBColor(1, 74, 129)		# #014A81 darker blue
#		fill.fore_color.rgb = RGBColor(28, 129, 185)	# light blue highlight for inactive
	else:
		fill.fore_color.rgb = RGBColor(191, 191, 191)	# grey for inactive bubbles
#		fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
#		fill.fore_color.brightness = -0.70

	shape.line.fill.background()	# no outline

	tf = shape.text_frame
	tf.margin_bottom = 0
	tf.margin_left = 0
	tf.margin_top = 0
	tf.margin_right = 0
	pg = tf.paragraphs[0]
	pg.line_spacing = line_spacing
	run = pg.add_run()
	run.text = text
	font = run.font
	font.name = 'Avenir Next LT Pro'
	font.size = Pt(fontsize)
	font.bold = True
	font.color.rgb = RGBColor(255, 255, 255)
#   if sing:
#		font.color.rgb = RGBColor(255, 255, 255)
#	else:
#		font.color.rgb = RGBColor(230, 230, 230)


def add_song_details(slide, displayBook, displaySong, meta, verses, chorus, coda, repeat):

	wp = get_song_bubbles(meta, verses, chorus, coda)	# tbd: handle codas

#	bwidth = (len(wp) * 0.48)
	bwidth = (len(wp) * 0.44) - 0.04
	bleft = 1.7 + (((10 - 1.7) - bwidth) / 2)
	btop = 3.725

	slotndx = 0
	for bubble in wp:
		add_bubble(slide, btop, bleft, slotndx, bubble["bubble"], bubble["sing"])
		slotndx = slotndx + 1



#
# add_texts
#

def add_texts(slide, language, texts, mapping, default_fonts):
	for ph in phlist[language]:
		v = get_placeholder(slide, mapping[ph])
		if ph in texts:
			text = texts[ph]
			if 'size' in text:
				size = text['size']
				set_placeholder_size(v, size[0], size[1], size[2], size[3])	# top, left, width, height)
			bold = text['bold'] if 'bold' in text else False
			fonts = text['fonts'] if 'fonts' in text else default_fonts
			font = fonts[1] if bold else fonts[0]
			features = text['features'] if 'features' in text else None
			if 'text' in text:
				v.text_frame.clear()
				run = v.text_frame.paragraphs[0].add_run()
				run.font.name = font[0]
				run.font.size = Pt(text['max_size'])
				run.text = text['text']
				if 'align' in text:
					v.text_frame.paragraphs[0].alignment = text['align']
				fitText(v.text_frame, font_family=font[0], max_size=text['max_size'], bold=bold, step_size=text['step_size'], file=assetRoot + font[1], features=features)
			if 'frame' in text:
				color = text['color'] if 'color' in text else None
				spacing = text['spacing'] if 'spacing' in text else 0
				fitTextRuns(v.text_frame, text['frame'], font_family=font[0], max_size=text['max_size'], bold=bold, step_size=text['step_size'], file=assetRoot + font[1], features=features, color=color, spacing=spacing)
		else:
			v._element.getparent().remove(v._element)

def get_leader(item):
	leader = None
	if 'leader' in item:
		leader = dict()
		leader[LAYOUT_TITLE_LEADER] = dict(text=item['leader'], max_size=20, step_size=2)
		# if item is a song, reposition leader box and alignment:
		if item['type'] == 'song' or item['type'] == 'song-title':
			leader[LAYOUT_TITLE_LEADER]['align'] = PP_ALIGN.CENTER
			leader[LAYOUT_TITLE_LEADER]['size'] = [5.6, 4.35, 3.0, 0.44]
	return leader

def set_logo(slide, ndx, logo):
	v = get_placeholder(slide, ndx)
	logo_png = assetRoot + logo
	set_placeholder_pic(v, logo_png, v.top.inches, v.left.inches, v.width.inches, v.height.inches)


def add_title_slide(outp, item, navitems, ndi, language, eng=None, esp=None, bgname="prayer-wheat", zoom=False, blur=False, delete=None, doLeader=True, goSlow=False):
	if 'style' in item:
		if item['style'] == "noslide":
			return

	if language == "bil":
		layout_ndx = LAYOUT_TITLE_BIL_SLOW if goSlow else LAYOUT_TITLE_BIL
	else:
		layout_ndx = LAYOUT_TITLE_ENG

	layout = outp.slide_masters[MASTER_TITLE].slide_layouts[layout_ndx]
	slide = outp.slides.add_slide(layout)

	mapping = ndx2ben	# first language slot
	# the three fixed items have the same mapping for both slide types

	background = get_background(bgname, item, zoom=zoom, blur=blur)

	# background image
	v = get_placeholder(slide, mapping[LAYOUT_TITLE_BACKGROUND])
	set_placeholder_pic(v, background, 0, 0, 10, 6.25)

	# clear background placeholders
	v = get_placeholder(slide, mapping[LAYOUT_TITLE_NAVBAR_BG])
	v.text = ' '
	v = get_placeholder(slide, mapping[LAYOUT_TITLE_CIRCLE_BG])
	v.text = ' '

	if language == "bil" or language == "eng":
		set_logo(slide, mapping[LAYOUT_TITLE_LOGO], "eh20-eng-432.png")
		add_texts(slide, "eng", eng, mapping, eng_fonts)
		mapping = ndx2bes	# second language slot

	if language == "bil" or language == "esp":
		if esp is None:
			esp = eng	# fall back to eng if no esp text
		set_logo(slide, mapping[LAYOUT_TITLE_LOGO], "eh20-esp-600.png")
		add_texts(slide, "esp", esp, mapping, esp_fonts)

	if doLeader:
		leader = get_leader(item)
		if leader:
			add_texts(slide, "ldr", leader, mapping, eng_fonts)
	else:
		leader = []
		add_texts(slide, "ldr", leader, mapping, eng_fonts)

	if delete is not None:
		for ph in delete:
			v = get_placeholder(slide, mapping[ph])
			v._element.getparent().remove(v._element)

	if 'style' in item:
		if item['style'] == "nosidebar":
			navitems = None

	create_sidebar(slide, language, navitems, ndi, layout_ndx)
	return slide






def add_song_title_slide(outp, language, item, navitems, ndi, verses=None, chorus=None, coda=None, repeat=None, bubble=None):
	displayBook = item['book']
	displaySong = int(item['song'])
	meta = item['meta']
	espm = item['esp']

	eng = None
	esp = None

	if language == "bil" or language == "eng" or espm is None:
		eng = dict()
		eng[LAYOUT_TITLE_TITLE] = dict(text=str(displaySong), max_size=60, step_size=6, bold=True, size=[0.675, 4.6, 2.5, 1])
		eng[LAYOUT_TITLE_DETAIL] = dict(text=meta['title'].upper(), max_size=54, step_size=6, bold=True, size=[1.675, 3.75, 4.2, 1.8])
		eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CSINGING...TO THE LORD\u201D\nCOLOSSIANS 3:16", max_size=24, step_size=2, size=[4.325, 4.1, 3.5, 0.7])
		eng[LAYOUT_TITLE_CREDITS] = dict(text='\n'.join(meta['credits'].splitlines()), max_size=12, step_size=2)
		if bubble:
			eng[LAYOUT_TITLE_CALLOUT] = dict(text=bubble.upper(), max_size=18, step_size=2, size=[0.25, 7.4, 3.0, 0.44], align=PP_ALIGN.LEFT)

	if language == "bil" or language == "esp":
		if espm:
			esp = dict()
			esp[LAYOUT_TITLE_TITLE] = dict(text=str(displaySong), fonts=eng_fonts, max_size=60, step_size=6, bold=True, size=[0.675, 4.6, 2.5, 1])
			esp[LAYOUT_TITLE_DETAIL] = dict(text=espm['title'].upper(), max_size=56, step_size=6, bold=True, size=[1.62, 3.75, 4.2, 1.8])
			esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CCANTANDO...AL SEÑOR\u201D\nCOLOSENSES 3:16", max_size=22, step_size=2, size=[4.280, 4.1, 3.5, 0.9])
			esp[LAYOUT_TITLE_CREDITS] = dict(text='\n'.join(espm['credits'].splitlines()), max_size=13, step_size=1)
			if bubble:
				esp[LAYOUT_TITLE_CALLOUT] = dict(text=espm['bubble'].upper(), max_size=18, step_size=2, size=[0.25, 7.4, 3.0, 0.44], align=PP_ALIGN.LEFT)

	doLeader = True if ('showLeader' in item) else False	# by default, do not show leader names on song titles

	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "song-beige", doLeader=doLeader)
	add_song_details(slide, displayBook, displaySong, meta, verses, chorus, coda, repeat)


#
# Old Medley Title Slide Support
#

def add_song_title(slide, song, pmap, bil):
	v = get_placeholder(slide, pmap[0])
	v.text = song['song']

	v = get_placeholder(slide, pmap[1])
	v.text = song['meta']['title']
	fitText(v.text_frame, font_family='Century Gothic', max_size=54, bold=False, step_size=6)

	v = get_placeholder(slide, pmap[2])
	v.text = '\n\n'.join(song['meta']['credits'].splitlines())
	fitText(v.text_frame, font_family='Century Gothic', max_size=16, bold=False, step_size=2)

	v = get_placeholder(slide, pmap[3])
	if song['book'] == 'eh':
		v._element.getparent().remove(v._element)
	else:
		book_png = assetRoot + song['book'] + "-rectangle.png"
		v.insert_picture(book_png)

	if bil == 1:
		if song['esp'] != None:
			esp = song['esp']
		else:
			esp = song['meta']

		v = get_placeholder(slide, pmap[4])
		v.text = esp['title']
		fitText(v.text_frame, font_family=None, max_size=54, bold=False, step_size=6, file=assetRoot + "AlegreyaSans-Medium.otf")

		v = get_placeholder(slide, pmap[5])
		v.text = '\n\n'.join(esp['credits'].splitlines())
		fitText(v.text_frame, font_family=None, max_size=16, bold=False, step_size=2, file=assetRoot + "AlegreyaSans-Light.otf")


def add_medley_title_slide(outp, item, navitems, ndi):
	s1 = item['songs'][0]
	s2 = item['songs'][1]

	# create and customize title slide
	if (s1['esp'] is None) and (s2['esp'] is None):
		layout = outp.slide_masters[4].slide_layouts[6]
		bil = 0
	else:
		layout = outp.slide_masters[4].slide_layouts[7]
		bil = 1
	slide = outp.slides.add_slide(layout)
#	dump_placeholders(slide)
#	create_navbar("med", slide, bil, navitems, ndi)

	add_song_title(slide, s1, [2,3,4,5,6,7], bil)
	add_song_title(slide, s2, [8,9,10,11,12,13], bil)




MASTER_SONG = 0

LAYOUT_SONG_BACKGROUND = 2
LAYOUT_SONG_TITLE_BAR = 3
LAYOUT_SONG_LOGO = 5
LAYOUT_SONG_NUMBER = 4
LAYOUT_SONG_TITLE = 6
LAYOUT_SONG_COPYRIGHT = 7
LAYOUT_SONG_NAVBAR = 8
LAYOUT_SONG_MUSIC = 9
LAYOUT_SONG_EH_LOGO = 10
LAYOUT_SONG_TITLE_ESP = 11
LAYOUT_SONG_COPYRIGHT_ESP = 12
LAYOUT_SONG_HIGHLIGHT = 13
LAYOUT_SONG_BUBBLE = 14
LAYOUT_SONG_BUBBLE_COUNT = 12
LAYOUT_SONG_ARROW = 26
LAYOUT_SONG_ARROW_COUNT = 4
LAYOUT_SONG_CODA = 30

arrow_count = 0
used_highlight = False
used_coda = False

def add_highlight(slide, top, left, height):
	global used_highlight
#   top = top - 0.04
#	left = left - 0.04
#	width = 0.48
	width = 0.40
	height = 1.36

	shape = get_placeholder(slide, LAYOUT_SONG_HIGHLIGHT)
	used_highlight = True

	fill = shape.fill
	fill.solid()
	fill.fore_color.rgb = RGBColor(28, 129, 185)	# light blue highlight for inactive
#	fill.background()
	shape.line.fill.background()	# no outline
#	shape.line.fill.solid()
#	shape.line.fill.fore_color.rgb = RGBColor(255, 255, 255)
	shape.text = ' '
	# must set size last, or else fill.solid() kills animation!
	set_placeholder_size(shape, top, left, width, height)

def add_arrow(slide, left, top):
	global arrow_count
	height = 0.4
	width = 0.22
	left = left + 0.09

	shape = get_placeholder(slide, LAYOUT_SONG_ARROW + arrow_count)
	arrow_count = arrow_count + 1

	fill = shape.fill
	fill.solid()
	fill.fore_color.rgb = RGBColor(255, 255, 255)
	shape.line.fill.background()	# no outline
	shape.text = ' '
	# must set size last, or else fill.solid() kills animation!
	set_placeholder_size(shape, top, left, width, height)

def add_circle(slide, slotndx, numslots, text, color, dw, dh):
	global used_coda
	diameter = 0.4
	slot_width = 0.48
	fontsize = 16
	line_spacing = 0.9

	slottop = 0.465 + dh
	top = slottop + (slotndx * slot_width)
	left = 0.25 + dw

	if color == "highlight":
		add_highlight(slide, top, left, slot_width * 3)

	shapes = slide.shapes
	shape = get_placeholder(slide, LAYOUT_SONG_BUBBLE + slotndx)

	if text == "arrow":
		add_arrow(slide, left, top)
		# delete corresponding bubble placeholder
		shape._element.getparent().remove(shape._element)
		return

	if text == "Coda":
		coda = get_placeholder(slide, LAYOUT_SONG_CODA)
		shape.text = ' '
		used_coda = True


	#1B255E -- eh dark blue
	#ED5959 -- eh red higlight

	fill = shape.fill
	fill.solid()
	if color == "blue" or color == "highlight":
		outline = fill.fore_color.rgb = RGBColor(28, 129, 185)	# light blue highlight for inactive
#		outline = fill.fore_color.rgb = RGBColor(27, 37, 94)	# eh dark blue
#		outline = fill.fore_color.rgb = RGBColor(28, 129, 185)	# light blue highlight
	elif color == "red":
		outline = fill.fore_color.rgb = RGBColor(237, 89, 89)	# eh red highlight
	elif color == "white":
		outline = fill.fore_color.rgb = RGBColor(255, 255, 255)	# when do we use this?
	else:
		outline = fill.fore_color.rgb = RGBColor(127, 127, 127)	# grey for inactive bubbles
#		fill.background()
#		outline = RGBColor(255, 255, 255)						# for "next verse" use white outline
#		outline = fill.fore_color.rgb = RGBColor(28, 129, 185)	# light blue highlight for inactive

	shape.line.fill.background()
#	shape.line.fill.solid()
#	shape.line.width = Pt(1.5)
#	shape.line.fill.fore_color.rgb = outline


	# must set size last, or else fill.solid() kills animation!
	set_placeholder_size(shape, top, left, diameter, diameter)

	if text == "Coda":
		set_placeholder_pic(coda, assetRoot + "coda-white.png", top, left, diameter, diameter)
		return
	if text == "Sanctus":
		text = "S"
	if text == "Amen":
		text = "A"
	tf = shape.text_frame
	tf.margin_bottom = 0
	tf.margin_left = 0
	tf.margin_top = 0
	tf.margin_right = 0
	pg = tf.paragraphs[0]
	pg.line_spacing = line_spacing
	run = pg.add_run()
	run.text = text
	font = run.font
	font.name = 'Avenir Next LT Pro'
	font.size = Pt(fontsize)
	font.bold = True
	if color == "white":
		font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
		font.color.brightness = -0.50
	else:
		font.color.rgb = RGBColor(255, 255, 255)


def invert_text(v, displaySong):
	fill = v.fill
	fill.solid()
	fill.fore_color.rgb = RGBColor(255, 255, 255)
	tf = v.text_frame
	pg = tf.paragraphs[0]
	run = pg.add_run()
	run.text = displaySong
	font = run.font
	font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
	font.color.brightness = -0.50

#
#
def add_verse_to_deck_tall(slide, displayBook, displaySong, basename, ndx, meta, esp, navinfo, bg_fn, invert, cdh):
	# top of tall column aligns with verse circles
	tall_top = cdh + 0.465

	# insert songbook logo and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_LOGO)
	book_png = assetRoot + displayBook + "-circle.png"
	top = tall_top
	left = 1.075
	set_placeholder_pic(v, book_png, top, left, v.width.inches, v.height.inches)

	# insert song number and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_NUMBER)
	if invert:
		invert_text(v, str(displaySong))
	else:
		v.text = str(displaySong)
	top = tall_top
	left = 0.125
	set_placeholder_size(v, top, left, v.width.inches, v.height.inches)
	fill = v.fill
	fill.background()

	tall_top  = tall_top + 0.5

	# insert song title and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_TITLE)
	top = tall_top
	left = 0.1
	width = 1.375
	height = v.height.inches
	set_placeholder_size(v, top, left, width, height)
	v.text = meta['title'].upper()
	fitText(v.text_frame, font_family=None, max_size=16, bold=True, step_size=1, file=assetRoot + "AvenirNextLTPro-Regular.ttf", features=('smcp', 'kern', 'liga', 'ordn', 'dlig'))

	tall_top = tall_top + height + 0.1

	# insert song title and adjust position (Spanish)
	v = get_placeholder(slide, LAYOUT_SONG_TITLE_ESP)
	if esp is None:
		v._element.getparent().remove(v._element)
	else:
		top = tall_top
		left = 0.1
		width = 1.375
		height = v.height.inches
		set_placeholder_size(v, top, left, width, height)
		v.text = esp['title'].upper()
		fitText(v.text_frame, font_family=None, max_size=18, bold=True, step_size=1, file=assetRoot + "AlegreyaSans-ExtraBold.ttf")

		tall_top = tall_top + height + 0.1


	# insert song copyright and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_COPYRIGHT)
	top = v.top.inches
	left = 0.1
	width = 1.375
	set_placeholder_size(v, top, left, width, v.height.inches)
	if meta['copyright'] == '':
		v._element.getparent().remove(v._element)
	else:
		v.text = meta['copyright']
		v.text_frame.margin_top = 0
		v.text_frame.margin_bottom = 0
		fitText(v.text_frame, font_family=None, max_size=8, bold=False, step_size=1, file=assetRoot + "AvenirNextLTPro-Regular.ttf")

	# insert song copyright and adjust position (Spanish)
	v = get_placeholder(slide, LAYOUT_SONG_COPYRIGHT_ESP)
	if esp is None:
		v._element.getparent().remove(v._element)
	else:
		if 'translation_copyright' in esp:
			top = v.top.inches
			left = 0.1
			width = 1.375
			set_placeholder_size(v, top, left, width, v.height.inches)
			v.text = "Traducción " + esp['translation_copyright']
		else:
			v._element.getparent().remove(v._element)

	# insert eh logo and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_EH_LOGO)
	top = v.top.inches
	left = 1.075
	set_placeholder_pic(v, assetRoot + "eh-circle.png", top, left, v.width.inches, v.height.inches)

	# set size & position of title bar background
	v = get_placeholder(slide, LAYOUT_SONG_TITLE_BAR)
	set_placeholder_size(v, 0, 0, 1.575, 6.25)
	v.text = ' '

#
def add_text_run(tf, text, fontname, size, bullet=False, bold=True):
	textrun = text
	if bullet:
		textrun = textrun + ' \u2022 '
	pg = tf.paragraphs[0]
	run = pg.add_run()
	run.text = textrun
	font = run.font
	font.name = fontname
	font.size = Pt(size)
	font.bold = bold
	font.color.rgb = RGBColor(255, 255, 255)

#
def add_wide_title(v, meta, esp):
	maxwidth = v.width
	engtitle = meta['title'].upper()
	if esp is None:
		add_text_run(v.text_frame, engtitle, 'Avenir Next LT Pro', 16, False)
	else:
		esptitle = esp['title'].upper()
		fontsize = 16
		engfontfile = assetRoot + "AvenirNextLTPro-Bold.otf"
		espfontfile = assetRoot + "AlegreyaSans-ExtraBold.ttf"

		while fontsize > 9:
			engw = _rendered_size(engtitle, fontsize, engfontfile)
			espw = _rendered_size(esptitle, fontsize + 2, espfontfile)
			if engw[0] + espw[0] < maxwidth:
				break
			else:
				fontsize = fontsize - 2

		add_text_run(v.text_frame, engtitle, 'Avenir Next LT Pro', fontsize, True)
		add_text_run(v.text_frame, esptitle, 'Alegreya Sans ExtraBold', fontsize + 2)




def add_verse_to_deck_wide(slide, displayBook, displaySong, basename, ndx, meta, esp, navinfo, bg_fn, invert):
	# insert songbook logo and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_LOGO)
	book_png = assetRoot + displayBook + "-circle.png"
	top = 5.75
	left = 0.8
	set_placeholder_pic(v, book_png, top, left, v.width.inches, v.height.inches)

	# insert song number and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_NUMBER)
	if invert:
		invert_text(v, str(displaySong))
	else:
		v.text = str(displaySong)
	top = 5.75
	left = -0.25
	set_placeholder_size(v, top, left, v.width.inches, v.height.inches)
	fill = v.fill
	fill.background()

	# insert song title and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_TITLE)
	top = 5.7
	left = 1.3
	width = 8.0	# v.width.inches - dw		# other half of width adjustment will be from esp title
	height = v.height.inches
	set_placeholder_size(v, top, left, width, height)

	add_wide_title(v, meta, esp)

	v = get_placeholder(slide, LAYOUT_SONG_TITLE_ESP)
	v._element.getparent().remove(v._element)


	# insert song copyright and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_COPYRIGHT)
	top = 6
	left = 1.4
	width = 8.0
	set_placeholder_size(v, top, left, width, v.height.inches)
	if meta['copyright'] == '' and esp is None:
		v._element.getparent().remove(v._element)
	else:
		espcopy = None
		if esp:
			if 'translation_copyright' in esp:
				espcopy = "Traducción " + esp['translation_copyright']
		if meta['copyright'] != '':
			add_text_run(v.text_frame, meta['copyright'], 'Avenir Next LT Pro', 8, espcopy is not None, bold=False)
		v.text_frame.margin_top = 0
		v.text_frame.margin_bottom = 0
		if espcopy:
			add_text_run(v.text_frame, espcopy, 'Alegreya Sans Medium', 9, bold=False)


	v = get_placeholder(slide, LAYOUT_SONG_COPYRIGHT_ESP)
	v._element.getparent().remove(v._element)


	# insert eh logo and adjust position
	v = get_placeholder(slide, LAYOUT_SONG_EH_LOGO)
	top = 5.75
	left = 9.4
	set_placeholder_pic(v, assetRoot + "eh-circle.png", top, left, v.width.inches, v.height.inches)

	# set size & position of title bar background
	v = get_placeholder(slide, LAYOUT_SONG_TITLE_BAR)
	set_placeholder_size(v, 6.25 - 0.6, 0, 10, 0.6)
	v.text = ' '




#
# add_verse_to_deck
#
def add_verse_to_deck(outp, displayBook, displaySong, basename, ndx, meta, esp, navinfo, slide, bg_fn, invert=False):
	global arrow_count
	global used_highlight
	global used_coda

	if slide:
		base_layout = 2
	else:
		base_layout = 0

	if esp is None:
		window = meta['window']
		orientation = meta['window_orientation']
	else:
		window = esp['window']
		orientation = esp['window_orientation']


	if orientation == 'wide':
		layout = outp.slide_masters[MASTER_SONG].slide_layouts[base_layout + 0]
		slide = outp.slides.add_slide(layout)
		cdw = (9.0 - window[2]) / 2.0
		cdh = (5.4 - window[3]) / 2.0 
		mw = 0
		add_verse_to_deck_wide(slide, displayBook, displaySong, basename, ndx, meta, esp, navinfo, bg_fn, invert)
	else:
		layout = outp.slide_masters[MASTER_SONG].slide_layouts[base_layout + 1]
		slide = outp.slides.add_slide(layout)
		cdw = (7.5 - window[2]) / 2.0 + 1.575
		cdh = (6.0 - window[3]) / 2.0 
		mw = 1.575
		add_verse_to_deck_tall(slide, displayBook, displaySong, basename, ndx, meta, esp, navinfo, bg_fn, invert, cdh)


	# insert music image and adjust size & position
	filename = basename + "-" + f"{ndx:02d}" + ".png"
	v = get_placeholder(slide, LAYOUT_SONG_MUSIC)
	set_placeholder_pic(v, filename, window[0], mw + window[1], window[2], window[3])

	# insert song background image
	v = get_placeholder(slide, LAYOUT_SONG_BACKGROUND)
	set_placeholder_pic(v, bg_fn, 0, 0, 10, 6.25)

	# add navigation circles
	ndy = 0
	arrow_count = 0
	used_highlight = False
	used_coda = False
	numslots = len(navinfo)
	for nav in navinfo:
		add_circle(slide, ndy, numslots, nav["text"], nav["color"], cdw, cdh)
		ndy = ndy + 1

	# remove any unused bubble placeholders
	while ndy < LAYOUT_SONG_BUBBLE_COUNT:
		v = get_placeholder(slide, LAYOUT_SONG_BUBBLE + ndy)
		v._element.getparent().remove(v._element)
		ndy = ndy + 1

	# remove any unused arrow placeholders
	while arrow_count < LAYOUT_SONG_ARROW_COUNT:
		v = get_placeholder(slide, LAYOUT_SONG_ARROW + arrow_count)
		v._element.getparent().remove(v._element)
		arrow_count = arrow_count + 1

	# remove the highlight placeholder, if not used
	if used_highlight is False:
		v = get_placeholder(slide, LAYOUT_SONG_HIGHLIGHT)
		v._element.getparent().remove(v._element)

	# remove the coda placeholder, if not used
	if used_coda is False:
		v = get_placeholder(slide, LAYOUT_SONG_CODA)
		v._element.getparent().remove(v._element)

	# remove navbar background
	v = get_placeholder(slide, LAYOUT_SONG_NAVBAR)
	v._element.getparent().remove(v._element)	# drw remove





def verse_iter(meta, verses):
	if verses:
		lastverse = 0
		for verse in verses:
			skip = (verse-1) != lastverse
			lastverse = verse
			yield str(verse), meta['verses'][str(verse)], skip
	else:
		lastverse = 0
		for verse, pngs in meta['verses'].items():
			skip = (int(verse)-1) != lastverse
			lastverse = int(verse)
			yield verse, pngs, skip

def add_song_to_deck(outp, item, verses=None, chorus=None, repeat=None, coda=True, sfv=True):
	displayBook = item['book']
	displaySong = int(item['song'])
	meta = item['meta']
	basename = item['basename']
	esp = item['esp']	# could be None?

	bg_fn = get_background("song-beige", item, blur=True)

	navinfo = get_navigation(meta, verses, chorus, coda)
	pprint.pprint(meta)
	if repeat:
		iterations = repeat + 1
	else:
		iterations = 1
	slide = sfv
	for iter in range(0, iterations):
		for verse, pngs, skip in verse_iter(meta, verses):
			for png in pngs:
				add_verse_to_deck(outp, displayBook, displaySong, basename, png, meta, esp, navinfo['nav'][png], slide, bg_fn)
				slide = False
			if verse in meta['chorus'].keys():
				if chorus and int(verse) not in chorus:
					continue
				for png in meta['chorus'][verse]:
					add_verse_to_deck(outp, displayBook, displaySong, basename, png, meta, esp, navinfo['nav'][png], slide, bg_fn)
					slide = False
		if coda and len(meta['codas']) > 0:
			for verse, pngs in meta['codas'].items():	# verse = "Coda" or "Amen" etc
				for png in pngs:
					add_verse_to_deck(outp, displayBook, displaySong, basename, png, meta, esp, navinfo['nav'][png], slide, bg_fn)
					slide = False


def add_medley_to_deck(outp, item, sfv=True):
	pprint.pprint(item)
#	navinfo = get_navigation(item['songs'][0]['meta'], [1,2], None, None)
	slide = sfv
	bg_fn = get_background("song-beige", item, blur=True)		# drw: tbd: revisit?
	for vx, verse in enumerate(item['nav']):
		si = item['songs'][verse[0]]
		if verse[1] == 'chorus':
			vi = si['meta']['chorus']['1']	# drw: assume first chorus always
		else:
			vi = si['meta']['verses'][verse[1]]

		nd = []
		for vy, vs in enumerate(item['nav']):
			if vx == vy:
				color = "blue"
			else:
				color = "gray" if vs[0] == 0 else "white"
			if vs[1] == 'chorus':
				vt = 'C'
			else:
				vt = vs[1]
			nd.append({ "text": vt, "color": color})

		for png in vi:
			add_verse_to_deck(outp, si['book'], int(si['song']), si['basename'], png, si['meta'], si['esp'], nd, slide, bg_fn, (verse[0] == 1))
			slide = False


#
# make_deck: emit a single song into a pptx
#

def set_esp(paths, language):
	esp = None
	basename = paths['engbase']
	espjson = paths['espbase'] + ".json"
	if os.path.exists(espjson):
		if language == 'esp' or language == 'bil':
			with open(espjson, 'r') as jsonfile:
				esp = json.load(jsonfile)
			pprint.pprint(esp)
			basename = paths['espbase']
	return esp, basename

def make_deck(book, number, language, outputfn):
	song, paths = get_song_paths_new(book, number)
	with open(paths['engbase'] + ".json", 'r') as jsonfile:
		meta = json.load(jsonfile)
	custom = paths['engbase'] + "-custom.json"
	if os.path.exists(custom):
		with open(custom, 'r') as jsonfile:
			meta.update(json.load(jsonfile))
	pprint.pprint(meta)

	esp, basename = set_esp(paths, language)

	verses = None
	chorus = None
	repeat = None
	coda = True

	outp = Presentation(assetRoot + "template-2020.pptx")

	item = dict()
	item['book'] = book
	item['song'] = str(number)
	item['meta'] = meta
	item['esp'] = esp
	item['basename'] = basename

	add_song_title_slide(outp, language, item, None, 0, verses, chorus, coda, repeat)
	add_song_to_deck(outp, item, verses, chorus, repeat, coda, True)
	outp.save(outputfn);



##############################################################################
##
## Worship Deck Generation
##
##############################################################################

#
# Template Masters & Layouts
#
# 0 - PowerPoint Default
# 8 - Standalone Slides
#   0 - Invitation
#   1 - Wed Announcements, push into 2
#   2 - Wed Please Join Us!, pushed from 1
#   3 - End Bumper (EH Logo & URL)
#   4 - Fade to Black
#

#
# 2020 Sidebar (navbar) logic
#

def set_highlight_size(w, item, htop, language):
	hwidths = { "eng" : { "welcome": 1.09, "reading": 0.98, "prayer": 0.77, "collection": 1.16, "supper": 1.50, "sermon": 0.91, "lesson": 0.83, "invitation": 1.05, "lyric": 0.68, "singing": 0.83 },
                "esp" : { "welcome": 1.22, "reading": 0.86, "prayer": 0.91, "collection": 0.93, "supper": 1.56, "sermon": 0.91, "lesson": 0.83, "invitation": 1.09, "lyric": 0.7,  "singing": 1.0  }}
	isSong = False
	if item[1] == 'song' or item[1] == 'song-title':
		isSong = True
		if item[3] == 'phss':
			deltaf = 0.2
		else:
			deltaf = 0.17
		delta = (3 - len(item[2])) * deltaf
		hwidth = 0.62 - delta
		hleft = 1.04 + delta
		# make room if needed for songbook icon
		if item[3] != 'pftl':
			hwidth = hwidth + 0.28
			hleft = hleft - 0.28
	elif item[1] in hwidths[language]:
		hwidth = hwidths[language][item[1]]
		hleft = 1.66 - hwidth
	else:
		hleft = w.left.inches
		hwidth = w.width.inches
	hheight = HIGHLIGHT_HEIGHT - 0.075
	htop = htop + (0.075 / 2)
#   if language == "esp" or (isSong and item[3] != 'pftl'):
	if language == "esp" or isSong:
		htop = htop - 0.02
	set_placeholder_size(w, htop, hleft, hwidth, hheight)

def add_sidebar_song(slide, item, bzero, ndx):
	if item[3] == 'phss':
		deltaf = 0.15
	else:
		deltaf = 0.2
	delta = (3 - len(item[2])) * deltaf
	diameter = Inches(0.4)
	left = Inches(0.72 + delta)
	top = Inches(bzero + (ndx * 0.375) + 0.04 - 0.07)
	shapes = slide.shapes
	shapes.add_picture(assetRoot + item[3] + "-xcircle-white.png", left, top, diameter)

def set_sidebar_item(slide, v, item, cycle, addp, bzero, ndx, language):
	p = v.text_frame.paragraphs[0] if not addp else v.text_frame.add_paragraph()
	if item[1] == 'song' or item[1] == 'song-title':
		if cycle == 0:
			if item[3] != 'pftl':
				add_sidebar_song(slide, item, bzero, ndx)
		number = item[2]

		p.line_spacing = Pt(18)
		p.space_before = Pt(1)
		p.space_after = Pt(8)

		run = p.add_run()
		run.text = number

		# always use eng font for numbers
		font = eng_fonts[1]	# bold

		run.font.name = font[0]
		run.font.size = Pt(18)
		run.font.bold = True
	else:
		if language == "esp":
			fonts = esp_fonts
		else:
			fonts = eng_fonts

		p.line_spacing = Pt(18)
		p.space_before = Pt(0)
		p.space_after = Pt(9)

		run = p.add_run()
		run.text = item[2]

		font = fonts[0]	# not bold
		run.font.name = font[0]
		run.font.size = Pt(16) if language == "eng" else Pt(18)
		run.font.bold = False
#		if len(item) > 3:
#			run.font.size = Pt(item[3])



#HIGHLIGHT_HEIGHT = 0.390625
#HIGHLIGHT_HEIGHT = 0.380625
#HIGHLIGHT_HEIGHT = 0.38
HIGHLIGHT_HEIGHT = 0.375


sidebar_phi = { LAYOUT_TITLE_ENG: [[LAYOUT_TITLE_ENG_NAVBAR, LAYOUT_TITLE_ENG_HIGHLIGHT]],
				LAYOUT_TITLE_BIL: [[LAYOUT_TITLE_BIL_ENG_NAVBAR, LAYOUT_TITLE_BIL_ENG_HIGHLIGHT],[LAYOUT_TITLE_BIL_ESP_NAVBAR, LAYOUT_TITLE_BIL_ESP_HIGHLIGHT] ],
				LAYOUT_TITLE_BIL_SLOW: [[LAYOUT_TITLE_BIL_ENG_NAVBAR, LAYOUT_TITLE_BIL_ENG_HIGHLIGHT],[LAYOUT_TITLE_BIL_ESP_NAVBAR, LAYOUT_TITLE_BIL_ESP_HIGHLIGHT] ],
				LAYOUT_WELCOME: [[LAYOUT_WELCOME_ENG_NAVBAR, LAYOUT_WELCOME_ENG_HIGHLIGHT]],
				LAYOUT_WELCOME_BIL: [[LAYOUT_WELCOME_ENG_NAVBAR, LAYOUT_WELCOME_ENG_HIGHLIGHT], [LAYOUT_WELCOME_ESP_NAVBAR, LAYOUT_WELCOME_ESP_HIGHLIGHT]] }

def set_sidebar(slide, items, myitem, cycle, phi, language):
	istart = 0
	bdelta = 0
	numitems = len(items)
	if numitems >= 14:
		if myitem > numitems - 14:
			istart = numitems - 14
		numitems = 14
	bheight = numitems * (HIGHLIGHT_HEIGHT)
	bzero = 0.08 + ((5.6 - 0.16) - bheight) / 2.0
	v = get_placeholder(slide, phi[0])
	v.text_frame.clear()
	addp = False

	sitems = items.copy()
	sitems.sort()

	for ndx in range(istart, numitems):
		set_sidebar_item(slide, v, sitems[ndx], cycle, addp, bzero, ndx, language)
		addp = True
		if myitem >= 0:
			highlight = (items[myitem][0] == (ndx + istart))
			if highlight:
				w = get_placeholder(slide, phi[1])
				w.text = ' '
				htop = bzero + (HIGHLIGHT_HEIGHT * ndx)
				set_highlight_size(w, sitems[ndx], htop, language)

#
def create_sidebar(slide, language, navitems, ndi, layout=LAYOUT_TITLE_ENG):
	phi = sidebar_phi[layout]
	if navitems is None:
		# remove background & higlight (for single songs)
		for phlist in phi:
			for ph in phlist:
				v = get_placeholder(slide, ph)
				v._element.getparent().remove(v._element)
	else:
		layer = 0
		if language == "bil" or language == "eng":
			set_sidebar(slide, navitems[0], ndi, layer, phi[layer], "eng")
			layer = 1
		if language == "bil" or language == "esp":
			set_sidebar(slide, navitems[1], ndi, layer, phi[layer], "esp")

		if ndi < 0:
			v = get_placeholder(slide, LAYOUT_WELCOME_ENG_HIGHLIGHT)
			v._element.getparent().remove(v._element)
			v = get_placeholder(slide, LAYOUT_WELCOME_ESP_HIGHLIGHT)
			v._element.getparent().remove(v._element)



#
# Slide Generation
#

def get_item(items, item):
	return items[item] if item in items else None

def add_song(outp, language, item, navitems, ndi, title=True, music=True):
	verses = get_item(item, 'verses')
	chorus = get_item(item, 'chorus')
	bubble = get_item(item, 'bubble')
	if item['esp'] != None:
		item['esp']['bubble'] = None
		if bubble != None:
			item['esp']['bubble'] = "Himno de Invitación"
	repeat = get_item(item, 'repeat')
	coda = get_item(item, 'coda')
	if coda is None:
		coda = False

	sfv = False		# only slide in first verse if we have title beforehand
	if title:
		add_song_title_slide(outp, language, item, navitems, ndi, verses, chorus, coda, repeat, bubble)
		sfv = True
	if music:
		add_song_to_deck(outp, item, verses, chorus, repeat, coda, sfv)


def add_song_entry(outp, language, item, navitems, ndi):
	eng = None
	esp = None

	if language == "bil" or language == "eng":
		eng = dict()
		eng[LAYOUT_TITLE_TITLE] = dict(text=' ', max_size=60, step_size=6, bold=True, size=[0.675, 4.6, 2.5, 1])
		eng[LAYOUT_TITLE_DETAIL] = dict(text=item['title'].upper(), max_size=54, step_size=6, bold=True, size=[1.675, 3.75, 4.2, 1.8])
		eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CSINGING...TO THE LORD\u201D\nCOLOSSIANS 3:16", max_size=24, step_size=2, size=[4.325, 4.1, 3.5, 0.7])
#		eng[LAYOUT_TITLE_CREDITS] = dict(text='\n'.join(meta['credits'].splitlines()), max_size=12, step_size=2)

	if language == "bil" or language == "esp":
		if 'título' in item:
			esp = dict()
			esp[LAYOUT_TITLE_TITLE] = dict(text=' ', fonts=eng_fonts, max_size=60, step_size=6, bold=True, size=[0.675, 4.6, 2.5, 1])
			esp[LAYOUT_TITLE_DETAIL] = dict(text=item['título'].upper(), max_size=56, step_size=6, bold=True, size=[1.62, 3.75, 4.2, 1.8])
			esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CCANTANDO...AL SEÑOR\u201D\nCOLOSENSES 3:16", max_size=22, step_size=2, size=[4.280, 4.1, 3.5, 0.9])
#			esp[LAYOUT_TITLE_CREDITS] = dict(text='\n'.join(espm['credits'].splitlines()), max_size=13, step_size=1)

	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "song-beige", doLeader=False)


#
# Medley Support
#

def add_medley(outp, item, navitems, ndi, title=True, music=True):
#	nav = get_item(item, 'nav')
	sfv = False		# only slide in first verse if we have title beforehand
	if title:
		add_medley_title_slide(outp, item, navitems, ndi)
		sfv = True
	if music:
		add_medley_to_deck(outp, item, sfv)




def add_prayer(outp, language, item, navitems, ndi):
	eng = None
	esp = None

	if language == "bil" or language == "eng":
		eng = dict()
		eng[LAYOUT_TITLE_DETAIL] = dict(text="PRAYER", max_size=80, step_size=6, bold=True, size=[1.675, 3.35, 5.0, 1.8])
		eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CDEVOTE YOURSELVES TO PRAYER\u201D", max_size=24, step_size=2, size=[3.325, 3.9, 3.9, 1.0])
		eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"COLOSSIANS 4:2", max_size=20, step_size=2, size=[4.225, 3.6, 4.52, 1.0])

	if language == "bil" or language == "esp":
		esp = dict()
		esp[LAYOUT_TITLE_DETAIL] = dict(text="ORACIÓN", max_size=80, step_size=6, bold=True, size=[1.590, 3.35, 5.0, 1.8])
		esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CPERSEVERAD EN LA ORACIÓN\u201D", max_size=28, step_size=2, size=[3.325, 3.9, 3.9, 1.0])
		esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"COLOSENSES 4:2", max_size=24, step_size=2, size=[4.225, 3.6, 4.52, 1.0])

	add_title_slide(outp, item, navitems, ndi, language, eng, esp, "prayer-wheat")


def add_sermon(outp, language, item, navitems, ndi):
	eng = None
	esp = None

	if item['type'] == 'lesson':
		engt = 'LESSON'
		espt = 'LECCIÓN'
	else:
		engt = 'SERMON'
		espt = 'SERMÓN'

	if language == "bil" or language == "eng":
		eng = dict()
		eng[LAYOUT_TITLE_DETAIL] = dict(text=engt, max_size=72, step_size=6, bold=True, size=[1.675, 3.35, 5.0, 1.8])
		eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CHE COMMANDED US TO PREACH TO THE PEOPLE\u201D", max_size=24, step_size=2, size=[3.33, 3.47, 4.78, 1.0])
		eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"ACTS 10:42", max_size=20, step_size=2, size=[4.225, 3.6, 4.52, 1.0])

	if language == "bil" or language == "esp":
		esp = dict()
		esp[LAYOUT_TITLE_DETAIL] = dict(text=espt, max_size=80, step_size=6, bold=True, size=[1.59, 3.35, 5.0, 1.8])
		esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CNOS MANDÓ QUE PREDICÁSEMOS AL PUEBLO\u201D", max_size=26, step_size=2, size=[3.29, 3.47, 4.78, 1.0])
		esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"HECHOS 10:42", max_size=24, step_size=2, size=[4.2, 3.6, 4.52, 1.0])

	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "sermon-notes")
	fade_to_black(outp, slide)





#
# Parsers to get readings from PPT files
#

def get_reading_text(file, ndx, esp):
	prs = Presentation(assetRoot + "readings-" + file + ".pptx")
	slide = prs.slides[ndx]
	if esp == 0:
		pr = get_placeholder(slide, 4)
		ps = get_placeholder(slide, 6)
	else:
		pr = get_placeholder(slide, 5)
		ps = get_placeholder(slide, 7)
	return [pr.text_frame, ps.text_frame]


def get_collection_reading(ndx, esp):
	return get_reading_text("collection", ndx, esp)

def get_supper_reading(ndx, esp):
	return get_reading_text("supper", ndx, esp)


def get_supper_reference(item, language):
	reference = " "
	if 'reading' in item:
		rndx = item['reading']
		if language == "eng":
			reading = get_supper_reading(rndx, 0)
		else:
			reading = get_supper_reading(rndx, 1)
		reference = reading[0].text
	return reference



def copy_paragraph(sp, dp, size, bold, color, spacing):
	dp.line_spacing = Pt(size + spacing)
	for srun in sp.runs:
		drun = dp.add_run()
		drun.text = srun.text
		if srun.font:
			drun.font.size = Pt(size)
			if srun.font.italic == True:
				drun.font.italic = True
			drun.font.bold = bold
			if color:
				drun.font.color.rgb = color


def copy_text_frame(src, dst, size=32, bold=False, color=None, spacing=0):
	p = src.paragraphs[0]
	copy_paragraph(p, dst.paragraphs[0], size, bold, color, spacing)
	for p in src.paragraphs[1:]:
		dp = dst.add_paragraph()
		copy_paragraph(p, dp, size, bold, color, spacing)


def fitTextRuns(dst, src, font_family, max_size, bold, step_size=12, file=None, features=None, color=None, spacing=0):
	# first pass: make our best guess on size
	dst.text = src.text
	size = fitText(dst, font_family, max_size, bold, step_size, file, features, spacing)
	print(dst.text, size)
	# second pass: copy each run (wtih formatting)
	dst.clear()
#	size = size - step_size		# just to be save
	copy_text_frame(src, dst, size, bold, color, spacing)



def fade_to_black(outp, slide):
	transition = parse_xml(
		'<p:transition spd="med" advTm="5000" %s>\n'
		"  <p:fade/>\n"
		"</p:transition>" % nsdecls("p")
		)
	slide.element._insert_transition(transition)
	# fade to black
	blayout = outp.slide_masters[MASTER_STATIC].slide_layouts[0]
	slide = outp.slides.add_slide(blayout)



#
def add_ls(outp, language, item, navitems, ndi):
	dointro = True
	if 'style' in item:
		if item['style'] == 'nointro':
			dointro = False

	# intro (talk)
	if dointro:
		eng = None
		esp = None

		if language == "bil" or language == "eng":
			eng = dict()
			eng[LAYOUT_TITLE_DETAIL] = dict(text="LORD\u2019S", max_size=72, step_size=6, bold=True, size=[1.455, 3.35, 5.0, 1.5])
			eng[LAYOUT_TITLE_REFERENCE] = dict(text="SUPPER", max_size=72, step_size=6, bold=True, size=[2.575, 3.35, 5.0, 2.4])
			eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CDO THIS IN REMEMBRANCE OF ME\u201D LUKE 22:19", max_size=24, step_size=2, size=[3.925, 3.67, 4.36, 1.0])

		if language == "bil" or language == "esp":
			esp = dict()
			esp[LAYOUT_TITLE_DETAIL] = dict(text="CENA DEL", max_size=72, step_size=6, bold=True, size=[1.300, 3.35, 5.0, 1.5])
			esp[LAYOUT_TITLE_REFERENCE] = dict(text="SEÑOR", max_size=72, step_size=6, bold=True, size=[2.575, 3.35, 5.0, 2.4])
			esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CHACED ESTO EN MEMORIA DE MÍ\u201D LUCAS 22:19", max_size=24, step_size=2, size=[3.925, 3.67, 4.36, 1.0])

		slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "supper-vial")
		fade_to_black(outp, slide)


	if 'reading' in item:
		eng = None
		esp = None

		rndx = item['reading']

		if language == "bil" or language == "eng":
			eng = dict()
			reading = get_supper_reading(rndx, 0)
			eng[LAYOUT_TITLE_TITLE] = dict(frame=reading[1], max_size=48, step_size=2, size=[0.2, 2.00, 7.75, 5.25], color=RGBColor(255, 255, 255), spacing=5)
			eng[LAYOUT_TITLE_CREDITS] = dict(frame=reading[0], max_size=18, step_size=2, size=[5.6, 7.2, 2.6, 0.4])

		if language == "bil" or language == "esp":
			esp = dict()
			reading = get_supper_reading(rndx, 1)
			esp[LAYOUT_TITLE_TITLE] = dict(frame=reading[1], max_size=48, step_size=2, size=[0.2, 2.00, 7.75, 5.25], color=RGBColor(255, 255, 255), spacing=5)
			esp[LAYOUT_TITLE_CREDITS] = dict(frame=reading[0], max_size=18, step_size=2, size=[5.6, 7.2, 2.6, 0.4])

		add_title_slide(outp, item, navitems, ndi, language, eng, esp, "supper-vial", zoom=True, delete=[LAYOUT_TITLE_CIRCLE_BG], doLeader=False, goSlow=True)




def add_giving(outp, language, item, navitems, ndi):
	eng = None
	esp = None

	if language == "bil" or language == "eng":
		eng = dict()
		eng[LAYOUT_TITLE_DETAIL] = dict(text="GIVING", max_size=80, step_size=6, bold=True, size=[1.205, 3.35, 5.0, 1.5])
		eng[LAYOUT_TITLE_QUOTE] = dict(text="THE COLLECTION FOR THE SAINTS", max_size=36, step_size=2, bold=True, size=[2.525, 3.53, 4.62, 1.3])
		eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"\u201CFOR GOD LOVES A CHEERFUL GIVER\u201D 2 CORINTHIANS 9:7", max_size=20, step_size=2, size=[3.925, 3.6, 4.52, 1.0])

	if language == "bil" or language == "esp":
		esp = dict()
		esp[LAYOUT_TITLE_DETAIL] = dict(text="OFRENDA", max_size=76, step_size=6, bold=True, size=[1.6, 3.31, 5.2, 1.5])
		esp[LAYOUT_TITLE_QUOTE] = dict(text="PARA LOS SANTOS", max_size=36, step_size=2, bold=True, size=[3.0, 3.53, 4.62, 0.85])
		esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"\u201CPORQUE DIOS AMA AL DADOR ALEGRE\u201D 2 CORINTIOS 9:7", max_size=24, step_size=2, size=[3.86, 3.6, 4.52, 1.0])

	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "giving-map")
	return slide

#
def add_zelle(outp, language, item, navitems, ndi):
	eng = dict()
	esp = dict()
	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "giving-options")

	# remove the background circle
	v = get_placeholder(slide, LAYOUT_TITLE_ENG_CIRCLE_BG)	# same for eng, esp, bil
	v._element.getparent().remove(v._element)

	return slide




#
def add_coll(outp, language, item, navitems, ndi):
	slide = None
	dointro = True
	dozelle = False
	if 'style' in item:
		if item['style'] == 'nointro':
			dointro = False
		elif item['style'] == 'zelle':
			dozelle = True

	# intro (talk)
	if dointro:
		if dozelle:
			slide = add_zelle(outp, language, item, navitems, ndi)
		else:
			slide = add_giving(outp, language, item, navitems, ndi)

	if 'reading' in item:
		if slide:
			fade_to_black(outp, slide)

		eng = None
		esp = None

		rndx = item['reading']

		if language == "bil" or language == "eng":
			eng = dict()
			reading = get_collection_reading(rndx, 0)
			eng[LAYOUT_TITLE_TITLE] = dict(frame=reading[1], max_size=48, step_size=2, size=[0.2, 2.00, 7.75, 5.25], color=RGBColor(255, 255, 255), spacing=5)
			eng[LAYOUT_TITLE_CREDITS] = dict(frame=reading[0], max_size=18, step_size=2, size=[5.6, 7.2, 2.6, 0.4])

		if language == "bil" or language == "esp":
			esp = dict()
			reading = get_collection_reading(rndx, 1)
			esp[LAYOUT_TITLE_TITLE] = dict(frame=reading[1], max_size=48, step_size=2, size=[0.2, 2.00, 7.75, 5.25], color=RGBColor(255, 255, 255), spacing=5)
			esp[LAYOUT_TITLE_CREDITS] = dict(frame=reading[0], max_size=18, step_size=2, size=[5.6, 7.2, 2.6, 0.4])

		add_title_slide(outp, item, navitems, ndi, language, eng, esp, "giving-map", zoom=True, delete=[LAYOUT_TITLE_CIRCLE_BG], doLeader=False, goSlow=True)






def add_invitation(outp, language, item, navitems, ndi):
	eng = None
	esp = None

	if language == "bil" or language == "eng":
		eng = dict()
		eng[LAYOUT_TITLE_DETAIL] = dict(text="INVITATION", max_size=60, step_size=6, bold=True, size=[1.675, 3.35, 5.0, 1.8])
		eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CCOME TO ME...AND I WILL GIVE YOU REST\u201D", max_size=24, step_size=2, size=[3.33, 3.47, 4.78, 1.0])
		eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"MATTHEW 11:28", max_size=20, step_size=2, size=[4.225, 3.6, 4.52, 1.0])

	if language == "bil" or language == "esp":
		esp = dict()
		esp[LAYOUT_TITLE_DETAIL] = dict(text="INVITACIÓN", max_size=68, step_size=6, bold=True, size=[1.59, 3.35, 5.0, 1.8])
		esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CVENID A MÍ...Y YO OS HARÉ DESCANSAR\u201D", max_size=26, step_size=2, size=[3.29, 3.47, 4.78, 1.0])
		esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"MATEO 11:28", max_size=24, step_size=2, size=[4.2, 3.6, 4.52, 1.0])

	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "invitation-road")
	fade_to_black(outp, slide)



def add_announcements(outp, item):
	layout = outp.slide_masters[MASTER_STATIC].slide_layouts[0]
	slide = outp.slides.add_slide(layout)
#	if 'style' in item and item['style'] == 'wed':
#		layout = outp.slide_masters[MASTER_STATIC].slide_layouts[2]
#		slide = outp.slides.add_slide(layout)

def add_fade(outp, item):
	layout = outp.slide_masters[MASTER_STATIC].slide_layouts[0]
	slide = outp.slides.add_slide(layout)

#def add_end_bumper(outp):
#	layout = outp.slide_masters[MASTER_STATIC].slide_layouts[3]
#	slide = outp.slides.add_slide(layout)


def add_scripture_reading(outp, language, item, navitems, ndi):
	eng = None
	esp = None

	if 'lang' in item:
		r_eng = item['lang'][0]
		r_esp = item['lang'][1]
	else:
		r_eng = item
		r_esp = None

	if language == "bil" or language == "eng":
		eng = dict()
		eng[LAYOUT_TITLE_DETAIL] = dict(text="READING", max_size=66, step_size=6, bold=True, size=[1.525, 3.35, 5.0, 1.0])
		eng[LAYOUT_TITLE_TITLE] = dict(text=r_eng['passage'].upper(), max_size=54, step_size=6, bold=False, size=[2.525, 3.5, 4.7, 1.2])
		if 'tag' in item:
			if item['tag'] == 'am':
				eng[LAYOUT_TITLE_QUOTE] = dict(text=u"NOT ASHAMED OF", bold=True, max_size=16, step_size=2, size=[4.385, 3.85, 4.0, 0.4])
				eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"GOD\u2019S COMMANDS", bold=True, max_size=16, step_size=2, size=[4.725, 4.35, 3.00, 0.5])
			else:
				eng[LAYOUT_TITLE_QUOTE] = dict(text=u"BOLD STATEMENTS OF", bold=True, max_size=16, step_size=2, size=[4.385, 3.85, 4.0, 0.4])
				eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"TRUST IN GOD", bold=True, max_size=16, step_size=2, size=[4.725, 4.35, 3.00, 0.5])
		else:
			eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CHAVE THIS LETTER READ TO ALL\u201D", max_size=16, step_size=2, size=[4.385, 3.85, 4.0, 0.4])
			eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"1 THESSALONIANS 5:27", max_size=16, step_size=2, size=[4.725, 4.35, 3.00, 0.5])
#		eng[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CHAVE THIS LETTER READ TO ALL\u201D", max_size=16, step_size=2, size=[4.385, 3.85, 4.0, 0.4])
		if 'pew' in r_eng:
			eng[LAYOUT_TITLE_CALLOUT] = dict(text="Pew Bible: Page " + r_eng['pew'], max_size=18, step_size=2, size=[3.77, 4.46, 2.87, 0.44])
#		eng[LAYOUT_TITLE_REFERENCE] = dict(text=u"1 THESSALONIANS 5:27", max_size=16, step_size=2, size=[4.725, 4.35, 3.00, 0.5])

	if language == "bil" or language == "esp":
		esp = dict()
		esp[LAYOUT_TITLE_DETAIL] = dict(text="LECTURA", max_size=72, step_size=6, bold=True, size=[1.34, 3.35, 5.0, 1.2])
		esp[LAYOUT_TITLE_TITLE] = dict(text=r_esp['passage'].upper(), max_size=66, step_size=6, bold=False, size=[2.525, 3.5, 4.7, 1.2])
		if 'tag' in item:
			if item['tag'] == 'am':
				esp[LAYOUT_TITLE_QUOTE] = dict(text=u"NO AVERGONZADOS DE LOS", bold=True, max_size=18, step_size=2, size=[4.385, 3.85, 4.0, 0.5])
				esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"MANDAMIENTOS DE DIOS", bold=True, max_size=18, step_size=2, size=[4.725, 4.35, 3.00, 0.5])
			else:
				esp[LAYOUT_TITLE_QUOTE] = dict(text=u"AFIRMACIONES VALIENTES DE", bold=True, max_size=18, step_size=2, size=[4.385, 3.85, 4.0, 0.5])
				esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"CONFIANZA EN DIOS", bold=True, max_size=18, step_size=2, size=[4.725, 4.35, 3.00, 0.5])
		else:
			esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CQUE ESTA CARTA SE LEA A TODOS\u201D", max_size=18, step_size=2, size=[4.385, 3.85, 4.0, 0.5])
			esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"1 TESALONICENSES 5:27", max_size=18, step_size=2, size=[4.725, 4.35, 3.00, 0.5])
#		esp[LAYOUT_TITLE_QUOTE] = dict(text=u"\u201CQUE ESTA CARTA SE LEA A TODOS\u201D", max_size=18, step_size=2, size=[4.385, 3.85, 4.0, 0.5])
#		esp[LAYOUT_TITLE_REFERENCE] = dict(text=u"1 TESALONICENSES 5:27", max_size=18, step_size=2, size=[4.725, 4.35, 3.00, 0.5])

	add_title_slide(outp, item, navitems, ndi, language, eng, esp, "podium")






##########################################################################
#
# Welcome Slide (old style)
#
##########################################################################

def getDisplayNumber(item):
	if item['book'] == "shs":
		return "S-" + item['song']
	elif item['book'] == "phss":
		return "PHSS-" + item['song']
	elif item['book'] == "eh":
		return "EH-" + item['song']
	return item['song']


tags_eng = dict(desc="desc", sermon="Sermon", lesson="Lesson", report="Report", title="title", reading="Scripture Reading", prayer="Prayer", announcements="Announcements", supper=u"Lord\u2019s Supper", collection="Collection from Members", invitation="Invitation", medley="Song Medley")
tags_esp = dict(desc="esp", sermon="Sermón", lesson="Lección", report="Reporte", title="título", reading="Lectura", prayer="Oración", announcements="Anuncios", supper="La Cena del Señor", collection="Ofrenda de los Miembros", invitation="Invitación", medley="Popurrí de canciones")


##########################################################################
#
# Welcome Slide (current style)
#
##########################################################################
class _Fonts(object):
    """
    A memoizing cache for ImageFont objects.
    """
    fonts = {}

    @classmethod
    def font(cls, font_path, point_size):
        if (font_path, point_size) not in cls.fonts:
            cls.fonts[(font_path, point_size)] = ImageFont.truetype(
                font_path, point_size
            )
        return cls.fonts[(font_path, point_size)]


def _rendered_size(text, point_size, font_file):
    """
    Return a (width, height) pair representing the size of *text* in English
    Metric Units (EMU) when rendered at *point_size* in the font defined in
    *font_file*.
    """
    emu_per_inch = 914400
    px_per_inch = 72.0

    font = _Fonts.font(font_file, point_size)

    # Calculate width and height from the bounding box
    left, top, right, bottom = font.getbbox(text)
    px_width = right - left
    px_height = bottom - top

#    px_width, px_height = font.getsize(text)

    emu_width = int(px_width / px_per_inch * emu_per_inch)
    emu_height = int(px_height / px_per_inch * emu_per_inch)

    return emu_width, emu_height

# add_welcome_item: add/format a single item of worship in the paragraph
#  item[0] is the text
#  item[1] is 0 for default color, 1 for song color (blue)
#  item[2] is leader


# add_welcome_order: generate text box with bulleted list of worship items
def add_worship_order(slide, v, order):
	tf = v.text_frame
	tf.clear()
	spacing = Pt(33.12)
	fontsize = Pt(22)
	fontbold = True if order[0][1] == 1 else False
	iy = 0
	p = tf.paragraphs[0]
	run = p.add_run()
	run.text = order[0][0]
	run.font.size = fontsize
	run.font.bold = fontbold
	for item in order[1:]:
		if iy < 13:
			p = tf.add_paragraph()
			iy = iy + 1
			run = p.add_run()
			run.text = item[0]
			run.font.size = fontsize
			fontbold = True if item[1] == 1 else False
			run.font.bold = fontbold


# add_welcome_order: generate text box with bulleted list of worship leaders
def add_worship_leaders(slide, v, order):
	tf = v.text_frame
	tf.clear()
	spacing = Pt(33.12)
	fontsize = Pt(22)
	iy = 0
	p = tf.paragraphs[0]
	run = p.add_run()
	run.text = order[0][2]
	run.font.size = fontsize
	for item in order[1:]:
		if iy < 13:
			p = tf.add_paragraph()
			iy = iy + 1
			run = p.add_run()
			run.text = item[2]
			run.font.size = fontsize


def parse_worship_item(order, item, language):
	tags = tags_eng if language == 'eng' else tags_esp
	leader = item['leader'] if 'leader' in item else " "
	if item['type'] == 'song' or item['type'] == 'song-music':
		if (language == 'esp') and (item['esp'] is not None):
			meta = item['esp']
		else:
			meta = item['meta']
		if 'short-title' in meta:
			title = meta['short-title']
		else:
			title = meta['title']
		order.append([title, 1, leader])
	elif item['type'] == 'lyric' or item['type'] == 'singing':
		if (language == 'esp') and ('título' in item):
			title = item['título']
		else:
			title = item['title']
		order.append([title, 1, leader])
	elif item['type'] == 'medley':
		desc = tags['medley'] + ": " + getDisplayNumber(item['songs'][0]) + " & " + getDisplayNumber(item['songs'][1])
		order.append([desc, 0, leader])
	elif item['type'] == 'sermon':
		if tags['title'] in item:
			desc = item[tags['title']]
		else:
			desc = tags['sermon']
		order.append([desc, 0, leader])
	elif item['type'] == 'lesson':
		if tags['title'] in item:
			desc = item[tags['title']]
		else:
			desc = tags['lesson']
		order.append([desc, 0, leader])
	elif item['type'] == 'report':
		if tags['title'] in item:
			desc = item[tags['title']]
		else:
			desc = tags['report']
		order.append([desc, 0, leader])
	elif item['type'] == 'invitation':
		if tags['title'] in item:
			desc = item[tags['title']]
		else:
			desc = tags['invitation']
		order.append([desc, 0, leader])
	elif item['type'] == 'prayer':
		desc = item[tags['desc']] if tags['desc'] in item else tags['prayer']
#		order.append([desc, 0])
		order.append([" ", 0, leader])
	elif item['type'] == 'reading':
		if 'lang' in item:
			ndx = 0 if language == "eng" else 1
			passage = item['lang'][ndx]['passage']
		else:
			passage = item['passage']
		order.append([passage, 0, leader])
	elif item['type'] == 'welcome':
		if tags['desc'] in item:
			desc = item[tags['desc']]
		else:
			desc = "Welcome"
		order.append([desc, 0, leader])
	elif item['type'] == 'announcements':
		order.append([tags['announcements'], 0, leader])
	elif item['type'] == 'ls-am-talk':
		if 'show' in item:
			if item['show'] == 1:
				order.append([u"Lord\u2019s Supper Talk", 0, leader])
	elif item['type'] == 'ls-pm' or item['type'] == 'ls-am':
		reference = get_supper_reference(item, language)
		order.append([reference, 0, leader])
	elif item['type'] == 'collection':
#		order.append([tags['collection'], 0])
		order.append([" ", 0, leader])


# add_visitor_card: add "visitor card" slide as sub-slide of welcome slide
def add_visitor_card(outp, language, item, navitems, ndi):
	eng = dict()
	esp = dict()
	slide = add_title_slide(outp, item, navitems, ndi, language, eng, esp, "visitor-background", doLeader=False)

	# remove the background circle
	v = get_placeholder(slide, LAYOUT_TITLE_ENG_CIRCLE_BG)	# same for eng, esp, bil
	v._element.getparent().remove(v._element)

	return slide


# add_welcome: add "welcome to service" slide displaying order of worship
def add_welcome(outp, language, worship, navitems, ndi):
	order = []
	orden = []
	for item in worship['items']:
		parse_worship_item(order, item, "eng")
		if language != "eng":
			parse_worship_item(orden, item, "esp")

	pprint.pprint(order)
	pprint.pprint(orden)

	if language == "bil":
		layout_ndx = LAYOUT_WELCOME_BIL
	else:
		layout_ndx = LAYOUT_WELCOME

	layout = outp.slide_masters[MASTER_TITLE].slide_layouts[layout_ndx]
	slide = outp.slides.add_slide(layout)

	# background image
	background = get_background("welcome-brown", item=worship['items'][0])
#	background = get_background("river-mountains")
	v = get_placeholder(slide, LAYOUT_WELCOME_BACKGROUND)
	set_placeholder_pic(v, background, 0, 0, 10, 6.25)

	# clear background placeholders
	v = get_placeholder(slide, LAYOUT_WELCOME_NAVBAR_BG)
	v.text = ' '

	# generate date string (TBD Spanish)
	if 'isodate' in worship:
		wdiso = datetime.fromisoformat(worship['isodate'])
		wd_eng = wdiso.strftime("%B ") + wdiso.strftime("%d, ").lstrip("0") + wdiso.strftime("%Y")
		wd_eng = wd_eng + "\n" + wdiso.strftime("%A") + " \u2022 " + wdiso.strftime("%I").lstrip("0") + wdiso.strftime(":%M") + wdiso.strftime("%p").lower()

		locale.setlocale(locale.LC_ALL, 'es_US')
		wd_esp = wdiso.strftime("%d ").lstrip("0") + wdiso.strftime("de %B")
		wd_esp = wd_esp + "\n" + wdiso.strftime("%A") + " \u2022 " + wdiso.strftime("%I").lstrip("0") + wdiso.strftime(":%M") + wdiso.strftime("%p").lower()
		locale.setlocale(locale.LC_ALL, 'en_US')
	else:
		wd_eng = worship['date']		# for legacy compatibility
		wd_esp = worship['date']		# for legacy compatibility


	# populate leaders (just one for both bil and single language)
	v = get_placeholder(slide, LAYOUT_WELCOME_LEADERS)
	add_worship_leaders(slide, v, order)


	mapping = [LAYOUT_WELCOME_ENG_LOGO, LAYOUT_WELCOME_ENG_DATE, LAYOUT_WELCOME_ENG_SILENCE, LAYOUT_WELCOME_ENG_ORDER]

	if language == "bil" or language == "eng":
		set_logo(slide, mapping[0], "eh20-eng-432.png")
		v = get_placeholder(slide, mapping[1])
		v.text = wd_eng
		v = get_placeholder(slide, mapping[2])
		v.text = "Please silence your phone"
		v = get_placeholder(slide, mapping[3])
		add_worship_order(slide, v, order)

		mapping = [LAYOUT_WELCOME_ESP_LOGO, LAYOUT_WELCOME_ESP_DATE, LAYOUT_WELCOME_ESP_SILENCE, LAYOUT_WELCOME_ESP_ORDER]

	if language == "bil" or language == "esp":
		set_logo(slide, mapping[0], "eh20-esp-600.png")
		v = get_placeholder(slide, mapping[1])
		v.text = wd_esp
		v = get_placeholder(slide, mapping[2])
		v.text = "Favor de apagar los celulares"
		v = get_placeholder(slide, mapping[3])
		add_worship_order(slide, v, orden)



#	if language != "eng":
#		if 'isodate' in worship:
#			locale.setlocale(locale.LC_ALL, 'es_US')
#			wd = wdiso.strftime("%B ") + wdiso.strftime("%d, ").lstrip("0") + wdiso.strftime("%Y")
#			locale.setlocale(locale.LC_ALL, 'en_US')
#		v = get_placeholder(slide, 4)
#		v.text = worship['servicio'] + "\n" + wd
#		v = get_placeholder(slide, 5)
#		add_worship_order(slide, v, orden)

	create_sidebar(slide, language, navitems, ndi, layout_ndx)

	# append the song leader name to the notes of this slide
	if 'leader' in worship:
		notes_slide = slide.notes_slide
		text_frame = notes_slide.notes_text_frame
		text_frame.text = "Song Leader: " + worship['leader']

	if worship["type"] == "Sun - AM":
		add_visitor_card(outp, language, worship['items'][0], navitems, ndi)


##########################################################################
#
# Worship Deck
#
##########################################################################

def get_navbar(worship, language):
	engitems = []
	espitems = []
	ndi = 0
	for item in worship['items']:
		if item['type'] == 'song':
			song = [ndi, "song", item["song"], item["book"]]
			engitems.append(song)
			espitems.append(song)
		elif item['type'] == 'lyric':
			engitems.append([ndi, "lyric", "Song"])
			espitems.append([ndi, "lyric", "Canto"])
		elif item['type'] == 'singing':
			engitems.append([ndi, "singing", "Singing"])
			espitems.append([ndi, "singing", "Cantando"])
		elif item['type'] == 'song-title':
			if 'skip' in item:
				skip = item['skip']
			else:
				skip = 2
			song = [ndi+skip, "song-title", item["song"], item["book"]]
			engitems.append(song)
			espitems.append(song)
			ndi = ndi - 1
		elif item['type'] == 'medley':
			engitems.append([ndi, "medley", "Medley"])
			espitems.append([ndi, "medley", "Popurrí"])
		elif item['type'] == 'prayer':
			engitems.append([ndi, "prayer", "Prayer"])
			espitems.append([ndi, "prayer", "Oración"])
		elif item['type'] == 'reading':
			engitems.append([ndi, "reading", "Reading"])
			espitems.append([ndi, "reading", "Lectura"])
		elif item['type'] == 'sermon':
			engitems.append([ndi, "sermon", "Sermon"])
			espitems.append([ndi, "sermon", "Sermón"])
		elif item['type'] == 'lesson':
			engitems.append([ndi, "lesson", "Lesson"])
			espitems.append([ndi, "lesson", "Lección"])
		elif item['type'] == 'report':
			engitems.append([ndi, "report", "Report"])
			espitems.append([ndi, "report", "Reporte"])
		elif item['type'] == 'invitation':
			engitems.append([ndi, "invitation", "Invitation"])
			espitems.append([ndi, "invitation", "Invitación", 11])
		elif item['type'] == 'welcome':
			engitems.append([ndi, "welcome", "Welcome"])
			espitems.append([ndi, "welcome", "Bienvenida", 12])
		elif item['type'] == 'ls-am' or item['type'] == 'ls-pm':
			engitems.append([ndi, "supper", u"Lord\u2019s Supper", 10.5])
			espitems.append([ndi, "supper", "Cena del Señor", 11])
		elif item['type'] == 'collection':
			engitems.append([ndi, "collection", "Collection", 11])
			espitems.append([ndi, "collection", "Ofrenda"])
		elif item['type'] == 'announcements':
			engitems.append([ndi, "announcements", "Closing", 10])
			espitems.append([ndi, "announcements", "Anuncios"])
		ndi = ndi + 1
	return [engitems, espitems]

def make_worship_deck(jsonfile):
	# set output filename from input
	outfn = os.path.splitext(jsonfile)[0] + ".pptx"

	with open(jsonfile, 'r') as file:
		worship = json.load(file)
	pprint.pprint(worship)

	language = worship['language'] if 'language' in worship else 'eng'

	# pre-fetch the meta data for any songs
	for item in worship['items']:
		if 'song' in item['type']:
			song, paths = get_song_paths_new(item['book'], int(item['song']))
			with open(paths['engbase'] + ".json", 'r') as jsonfile:
				item['meta'] = json.load(jsonfile)
				custom = paths['engbase'] + "-custom.json"
				if os.path.exists(custom):
					with open(custom, 'r') as jsonfile:
						item['meta'].update(json.load(jsonfile))
			item['esp'], item['basename'] = set_esp(paths, language)
		elif 'medley' in item['type']:
			for songi in item['songs']:
				song, paths = get_song_paths_new(songi['book'], int(songi['song']))
				with open(paths['engbase'] + ".json", 'r') as jsonfile:
					songi['meta'] = json.load(jsonfile)
					custom = paths['engbase'] + "-custom.json"
					if os.path.exists(custom):
						with open(custom, 'r') as jsonfile:
							songi['meta'].update(json.load(jsonfile))
				songi['esp'], songi['basename'] = set_esp(paths, language)

	outp = Presentation(assetRoot + "template-2020.pptx")

	if 'nonav' in worship:
		navitems = None
	else:
		navitems = get_navbar(worship, language)
		print(navitems)

	# hook to add welcome page without a welcome item in the order of worship
	#   the -1 makes sure we do not highlight anything in the sidebar
	if 'welcome' in worship:
		add_welcome(outp, language, worship, navitems, -1)

	ndi = 0
	for item in worship['items']:
#drw		print("myitem", ndi)
		if item['type'] == 'song':
			add_song(outp, language, item, navitems, ndi)
		elif item['type'] == 'song-title':
			add_song(outp, language, item, navitems, ndi, music=False)
		elif item['type'] == 'song-music':
			add_song(outp, language, item, navitems, ndi, title=False)
			ndi = ndi - 1
		elif item['type'] == 'lyric':
			add_song_entry(outp, language, item, navitems, ndi)
		elif item['type'] == 'medley':
			add_medley(outp, item, navitems, ndi)
		elif item['type'] == 'prayer':
			add_prayer(outp, language, item, navitems, ndi)
		elif item['type'] == 'reading':
			add_scripture_reading(outp, language, item, navitems, ndi)
		elif item['type'] == 'sermon' or item['type'] == 'lesson' or item['type'] == 'report':
			add_sermon(outp, language, item, navitems, ndi)
		elif item['type'] == 'invitation':
			add_invitation(outp, language, item, navitems, ndi)
		elif item['type'] == 'welcome':
#			add_welcome(outp, item, worship, language)
			add_welcome(outp, language, worship, navitems, ndi)
		elif item['type'] == 'ls-am':
			add_ls(outp, language, item, navitems, ndi)
		elif item['type'] == 'collection':
			add_coll(outp, language, item, navitems, ndi)
		elif item['type'] == 'announcements':
			add_announcements(outp, item)

		if 'fade' in item:
			add_fade(outp, item)
		ndi = ndi + 1

	# setup for countdown timer
	# diff = Now - #10/17/2020 2:15:00 PM#
	if 'isodate' in worship:
		wdiso = datetime.fromisoformat(worship['isodate'])
		wd = wdiso.strftime("#%m/%d/%Y %H:%M:%S#")
		print(wd)

#	add_end_bumper(outp)
	outp.save(outfn);


##############################################################################
##
## Bilingual (Spanish) Slide Processing
##
##############################################################################

def hint(hints, fndx, staff, value):
	slide = str(fndx)
	hint = None
	if hints is not None:
		if value in hints:
			hint = hints[value]
		if slide in hints:
			if value in hints[slide]:
				hint = hints[slide][value]
			if staff in hints[slide]:
				if value in hints[slide][staff]:
					hint = hints[slide][staff][value]
	return hint


# return a list of black pixel ranges on the given row
#  (skip over the first pw pixels - the left staff)
def get_blobs(pixels, imgw, mode, width, pw, y, hintpixel):
	row = pixels[(y * imgw):((y+1) * imgw)]
	blobs = []
	blob = None
#	for ndx in range(pw, len(row)):
	for ndx in range(pw, pw + width):
#		if ((mode == '1' or mode == 'L') and row[ndx] != 255) or ((mode == 'RGBA' or mode == 'RGB') and row[ndx] != (255,255,255,255)) or ((mode == 'P') and (row[ndx] != 1)):
		if ((mode == '1' or mode == 'L') and row[ndx] != 255) or (mode == 'RGBA' and row[ndx] != (255,255,255,255)) or ((mode == 'RGB') and row[ndx] != (255,255,255)) or ((mode == 'P') and (row[ndx] != hintpixel)):
			if blob is None:
				blob = [ ndx, ndx ]
			else:
				blob[1] = ndx
		else:
			if blob:
				blobs.append(blob)
				blob = None
	return blobs

def blobs_overlap(blob, lastblobs):
	overlap = False
	for lb in lastblobs:
		if blob[0] > lb[1]:
			continue
		if blob[1] < lb[0]:
			continue
		overlap = True
		break
	return overlap

def break_mask(img, imgw, width, breakline, pw, hintpixel):
	print("**** break_mask:", imgw, width, breakline, pw, hintpixel)
	pixels = list(img.getdata())
	mask = None
	# start at breakline-1 -- the last line we will copy
	lastblobs = get_blobs(pixels, imgw, img.mode, width, pw, breakline-1, hintpixel)

#	pprint.pprint(lastblobs)

	# move through one line at a time, finding overlapping blobs on each
	#   successive line
	mask_y = breakline
	mask_spec = dict()
	while True:
		nextblobs = []
		orphans = []
		blobs = get_blobs(pixels, imgw, img.mode, width, pw, mask_y, hintpixel)
		if len(blobs) == 0:
			break
		for blob in blobs:
			if blobs_overlap(blob, lastblobs):
				nextblobs.append(blob)
			else:
				orphans.append(blob)

		if len(nextblobs) > 0:
			mask_spec[mask_y] = [ nextblobs, orphans ]
			lastblobs = nextblobs
			mask_y = mask_y + 1
		else:
			break

	mask_h = mask_y - breakline

	# go back upwards toward the breakline to see if any orphans
	#   connect to the blobs below them
	if mask_h > 1:
		for y in range(mask_y - 2, breakline, -1):
			for blob in mask_spec[y][1]:
				if blobs_overlap(blob, mask_spec[y + 1][0]):
					mask_spec[y][0].append(blob)

#	pprint.pprint(mask_spec)

	mask_0 = 0
	mask_1 = 255

	if mask_h > 0:
		mask = PIL.Image.new("L", (width, mask_h), color=mask_0)
		for y, row in mask_spec.items():
			for blob in row[0]:
				for x in range(blob[0], blob[1]+1):
					mask.putpixel((x - pw,y-breakline), mask_1)

	return mask, mask_h

def filter_words(tsv, minconf, minlen):
	words = []
	rawcount = len(tsv['level'])
	for ndx in range(0, rawcount):
		if int(tsv['conf'][ndx]) > minconf and int(tsv['top'][ndx]) > 0 and len(tsv['text'][ndx]) >= minlen:
			words.append({'top' : int(tsv['top'][ndx]), 'height': int(tsv['height'][ndx]), 'text': tsv['text'][ndx]})
	return words

def expand_staff(hints, fndx, img, pixels, width, egbdf, topstaff, hintpixel):
	slide = str(fndx)						# for parsing hints
	staff = "top" if topstaff else "bot"	# for parsing hints

	# use the first line after the top line of the staff for padding
#	padline = egbdf[0]['top'] + egbdf[0]['height']
#	padline = 236	# drw: eh-107
	padline = hint(hints, fndx, staff, "padline")
	if padline is None:
		padline = egbdf[0]['top'] + egbdf[0]['height']
	print("padline", padline)
	row = pixels[(padline * width):(padline * width)+width]
	pw = row.index(0)
	pw = pw + row[pw:].index(255)
	pw = pw + row[pw:].index(0)

	# most pftls have two parts to staff; RJs have a solid single
	if hint(hints, fndx, staff, "singleleft") is None:
		pw = pw + row[pw:].index(255)
#	else:	# drw: why minus 2?
#		pw = pw - 2

	# to handle "double lines" on the left
	if hint(hints, fndx, staff, "doubleleft") is not None:
		print("doubleleft")
		pw = pw + row[pw:].index(0)
		pw = pw + row[pw:].index(255)

	pwdelta = hint(hints, fndx, staff, "pwdelta")
	if pwdelta is not None:
		pw = pw + pwdelta
#	if pw > 680:
#		pw = pw - 78	# drw: make an option

	top = egbdf[4]['top'] + egbdf[4]['height']
	bot = egbdf[5]['top'] 
	print("staff crop", 0, top, width, bot)
	lyric = img.crop(box=(0, top, width, bot))
#	lyric.save("lyric.png")
	# Get verbose data including boxes, confidences, line and page numbers
	tsv = pytesseract.image_to_data(lyric, output_type=pytesseract.Output.DICT, config='--dpi '+str(width/10) + ' --psm 12')
	print(tsv)
	words = filter_words(tsv, 80, 3)
#	words = filter_words(tsv, 60, 3)	# pftl-707
	if len(words) == 0:
		words = filter_words(tsv, 60, 3)
		if len(words) == 0:
			words = filter_words(tsv, 40, 3)
	pprint.pprint(words)
	bots = []
	toth = 0
	for word in words:
		# drw: patch for pftl-707
#		if word['height'] > 100:
#			word['height'] = 80
		bots.append(word['top'] + word['height'])
		toth = toth + word['height']
	bots.sort()

	if len(words) != 0:
		avgh = toth / len(words)
		print(bots)
		print(avgh)
		bottom = bots
		if bots[-1] - bots[0] > avgh:
			height = hint(hints, slide, staff, "height") or int((avgh * 2) * 1.6)	# allow space for two lines of lyrics
#			breakbottom = bots[0] + int((bots[-1] - bots[0]) / 2)
			for ndx in range(0, len(bots)-1):
				if (bots[ndx] + (avgh/2)) < bots[ndx + 1]:
					bottom = bots[ndx+1:]
					break
#			breakbottom = mode(bottom)
#			breakbottom = max(set(bottom), key = bottom.count)
			breakbottom = bottom[0]
#			breakbottom = breakbottom + 60	# drw; fudge factor for 368
#			height = int(height * 1.2)	# for Exalted
#			if not topstaff:
#				breakbottom = bottom[0] - 20	# drw: for pftl-621
#			if topstaff:
#				breakbottom = bottom[0] - 20	# drw: for pftl-621
		else:
			height = hint(hints, slide, staff, "height") or int(avgh * 1.6)
			breakbottom = bottom[0]
#			if img.size[1] > 2500:
#				breakbottom = bottom[0] + 40	# drw: eh-101
#			height = int(height * 2.1)	# for Exalted
#			breakbottom = breakbottom - 10	# drw; fudge factor for 368
#			breakbottom = breakbottom + 60	# drw; fudge factor for 368
		print(bottom)
	else:
		breakbottom = int(6 * (egbdf[5]['top'] - egbdf[4]['top']) / 10)
		height = hint(hints, slide, staff, "height") or int((egbdf[5]['top'] - egbdf[4]['top']) * 0.25)

	bdelta = hint(hints, slide, staff, "bdelta") or 0

	breakline = egbdf[4]['top'] + egbdf[4]['height'] + breakbottom + bdelta
	print(breakbottom, breakline)

	coldef = hint(hints, slide, staff, "coldef")
	columns = []

	print(coldef)
	if coldef is None:
		mask, mask_h = break_mask(img, width, width - pw, breakline, pw, hintpixel)
		columns.append({'insert': breakline, 'height': height, 'srcwidth': pw, 'padline': padline, 'mask': mask, "mask_h": mask_h, "x": pw, "width": width - pw})
	else:
		# first column
		mask, mask_h = break_mask(img, width, coldef[0]['width'] - pw, breakline + coldef[0]['delta'], pw, hintpixel)
		columns.append({'insert': breakline + coldef[0]['delta'], 'height': height, 'srcwidth': pw, 'padline': padline, 'mask': mask, "mask_h": mask_h, "x": pw, "width": coldef[0]['width'] - pw})

		cpw = coldef[0]['width']
		for ndx, col in enumerate(coldef[1:]):
			mask, mask_h = break_mask(img, width, col['width'], breakline + col['delta'], cpw, hintpixel)
			columns.append({'insert': breakline + col['delta'], 'height': height, 'mask': mask, "mask_h": mask_h, "x": cpw, "width": col['width'], "srcwidth": 0, "padline": padline})
			cpw = cpw + col['width']

	return columns


def addop_copy(ops, left, top, width, height, dy):
	ops.append({ "op": "copy", "top": top, "height": height, "left": left, "width": width, "dy": dy})
	return dy + height

def addop_mask(ops, sp, inverse, dy):
	ops.append({ "op": "mask", "mask": sp['mask'], "top": sp['insert'], "height": sp['mask_h'], "left": sp['x'], "width": sp['width'], "inverse": inverse, "srcwidth": sp['srcwidth'], "srcline": sp['padline'], "dy": dy})
	return dy + sp['mask_h']

def addop_split(ops, sp, dy):
	padh = sp['height']
	if sp['mask'] is not None:
		dy = addop_mask(ops, sp, 0, dy)
		padh = padh - sp['mask_h']
	ops.append({ "op": "pad", "srcline": sp['padline'], "srcwidth": sp['srcwidth'], "height": padh, "left": sp['x'], "width": sp['width'], "dy": dy})
	dy = dy + padh
	if sp['mask'] is not None:
		dy = addop_mask(ops, sp, 1, dy)
	return dy


def staff2ops(ops, spec, top, bot, dy_first):
	for sp in spec:
		sx = sp['x']
		sw = sp['width']
		if 'srcwidth' in sp:
			sx = sx - sp['srcwidth']
			sw = sw + sp['srcwidth']
		dy = addop_copy(ops, sx, top, sw, sp['insert'] - top, dy_first)
		dy = addop_split(ops, sp, dy)
		dy = addop_copy(ops, sx, sp['insert'] + sp['mask_h'], sw, bot - sp['insert'] - sp['mask_h'] + 1, dy)
		lastdy = dy
	return lastdy


def find_subtitle_space(picture, fndx, nfiles, hints):
	oimg = PIL.Image.open(picture)
	print(oimg.mode, oimg.size)
	if oimg.mode != "1":
		img = oimg.convert(mode="1")
	else:
		img = oimg
	print(img.mode, img.size)

	staffw = hint(hints, fndx, None, "staffw")
	if staffw is None:
		staffw = 3

	pixels = list(img.getdata())
	width, height = img.size
	top = -1
	bot = height + 1
	left = width + 1
	right = -1
	egbdf = []
	for ny in range(height):
		row = pixels[(ny * width):(ny * width)+width]
		np = row.count(255)		# white pixel count
		if np != width:
			if top == -1:
				top = ny
			bot = ny
#			if (((fndx == nfiles) and (np < 2000)) or ((fndx != nfiles) and (np < width / 3))):
#			if np < width / 2:	# drw: for eh-106, 585
			if np < width / staffw:
				if len(egbdf) == 0:
					egbdf.append({"top": ny, "white": [np], "height": 1})
				elif egbdf[-1]['top'] + egbdf[-1]['height'] == ny:
					egbdf[-1]['height'] = egbdf[-1]['height'] + 1   			  #		
					egbdf[-1]['white'].append(np)
				else:
					egbdf.append({"top": ny, "white": [np], "height": 1})
			tl = row.index(0)
			if tl < left:
				left = tl
			rt = width - row[::-1].index(0)
			if rt > right:
				right = rt
	print(top, bot, left, right)

	ops = []

	hp = hint(hints, fndx, None, "pixel")
	if hp is None:
		hp = 1

	if hint(hints, fndx, None, "passthru"):
		print("passthru")
		dy = addop_copy(ops, 0, top, width, height - top, 0)
	elif len(egbdf) == 20:
		spec1 = expand_staff(hints, fndx, oimg, pixels, width, egbdf[0:10], True, hp)
		spec2 = expand_staff(hints, fndx, oimg, pixels, width, egbdf[10:20], False, hp)
		pprint.pprint(spec1)
		pprint.pprint(spec2)

		border = egbdf[10]['top'] 
		dy = staff2ops(ops, spec1, top, border, 0)
		dy = staff2ops(ops, spec2, border, bot, dy)

	elif len(egbdf) == 10:
		spec = expand_staff(hints, fndx, oimg, pixels, width, egbdf, True, hp)
		dy = staff2ops(ops, spec, top, bot, 0)
	else:
		print("error: incorrect number of staff lines", len(egbdf), fndx)

	return ops, dy

def do_pad(op, src, dst):
#	print("do_pad", dst.size)
#	print(op)
	if op['srcwidth'] != 0:
		for py in range(0, op['height']):
			ty = op['dy'] + py
			if ty < dst.size[1]:
				for x in range(0, op['srcwidth']):
					pixel = src.getpixel((x, op['srcline']))
					dst.putpixel((x, ty), pixel)

def add_subtitle_space(picture, fndx, nfiles, hints):
	ops, height = find_subtitle_space(picture, fndx, nfiles, hints)

	pprint.pprint(ops)
	print("height", height)

	src = PIL.Image.open(picture)
	pprint.pprint(src.info)
	width, oheight = src.size
	if src.mode == 'P':
		white = hint(hints, fndx, None, "pixel")
		if white is None:	# need to allow white to be zero
			white = 1	# drw: or 0 (in some cases)
	else:
		white = PIL.ImageColor.getcolor("white", src.mode)
	dst = PIL.Image.new(src.mode, (width, height), color=white)

	if src.mode == 'P':
		palette = src.getpalette()
#		pprint.pprint(palette)
#		for ndx in range(0,256*3,3):
#			if palette[ndx] != 0 and palette[ndx+1] != 0 and palette[ndx+2] != 0:
#				palette[ndx] = 255
#				palette[ndx+1] = 255
#				palette[ndx+2] = 255
#			else:
#				palette[ndx] = 0
#				palette[ndx+1] = 0
#				palette[ndx+2] = 0
#		pprint.pprint(palette)
		dst.putpalette(palette)

	for op in ops:
		if op['op'] == 'copy':
			if op['height'] > 0:
				strip = src.crop(box=(op['left'], op['top'], op['left'] + op['width'], op['top'] + op['height']))
				dst.paste(strip, box=(op['left'], op['dy'], op['left'] + op['width'], op['dy'] + op['height']))
		elif op['op'] == 'mask':
			strip = src.crop(box=(op['left'], op['top'], op['left'] + op['width'], op['top'] + op['height']))
			if op['inverse']:
				mask = PIL.ImageOps.invert(op['mask'])
			else:
				mask = op['mask']
			dst.paste(strip, box=(op['left'], op['dy'], op['left'] + op['width'], op['dy'] + op['height']), mask=mask)
			do_pad(op, src, dst)	# fill in the left staff as well, like when we pad
		elif op['op'] == 'pad':
			do_pad(op, src, dst)

	return dst


def make_esp_blank(book, number, target):
	song, pathname, basename, rawname = get_song_paths(book, number)
	with open(basename + ".json", 'r') as jsonfile:
		meta = json.load(jsonfile)
	pprint.pprint(meta)

	# Fixed directories (one file per son)
	# esp/pftl/eng/book-xyz-eng.pptx (tool generated deck with expanded space for subtitles)
	# esp/pftl/bil/book-xyz-bil.pptx (deck with embedded text boxes for subtitles)

	# Per song directories
	# esp/pftl/bil/xyz/*.png, *.txt (intermediate output from exporting from pptx file)
	# esp/pftl/xyz/ -- json and pngs

	# Create a directory to hold our output
	bilpath = '/'.join(['ehsf', 'esp', book, 'bil', song])
	os.makedirs(bilpath, exist_ok=True)
	esppath = '/'.join(['ehsf', 'esp', book, song])
	os.makedirs(esppath, exist_ok=True)

	# Generate output file names and hint fine name
	ofbase = "ehsf/esp/" + book + "/eng/" + book + "-" + song
	ofpptx = ofbase + "-eng.pptx"
	ofhint = ofbase + "-hints.json"

	# Read in optional "hints" if they exist
	hints = None
	if os.path.exists(ofhint):
		with open(ofhint, 'r') as hintfile:
			hints = json.load(hintfile)
		pprint.pprint(hints)

	outp = Presentation(assetRoot + "template-spanish.pptx")
	blank_slide_layout = outp.slide_layouts[6]
	swidth = outp.slide_width

	# add metadata slide
	layout = outp.slide_masters[1].slide_layouts[0]
	slide = outp.slides.add_slide(layout)
	for v in slide.placeholders:
		if v.shape_id == 2:
			v.text = meta['title']
		elif v.shape_id == 3:
			v.text = meta['credits']

	# take images from unprocessed, original pptx (stored in /raw)
	files = glob.glob(rawname + "*.png")
	files.sort()
	pngs = []
	nfiles = len(files)
	for ndx, file in enumerate(files, 1):
		if target is not None:
			if ndx != int(target):
				continue
		with open(file, 'rb') as f:
			pngs.append(add_subtitle_space(f, ndx, nfiles, hints))

	minar = 4.0/3.0					# default aspect ratio
	for ndx, png in enumerate(pngs, 1):
		thisar = png.size[0] / png.size[1]
		print(png.mode, png.size, thisar)
		if thisar < minar:
			minar = thisar
	print("minar", minar)
	outp.slide_height = Inches(10.0/minar)

	for ndx, png in enumerate(pngs, 1):
		xw = png.size[0]
		xh = int((png.size[0]/minar) + 0.5)
		dpi = xw/10
		oimg = io.BytesIO()
		png.save(oimg, format="png", dpi=(dpi,dpi))

		oslide = outp.slides.add_slide(blank_slide_layout)
		pic = oslide.shapes.add_picture(oimg, 0, 0, width=swidth)

		fn = song + "/" + book + "-" + song + "-" + f"{ndx:03d}" + ".png" 
		notes_slide = oslide.notes_slide
		text_frame = notes_slide.notes_text_frame
		text_frame.text = str(xw) + "x" + str(xh) + "x" + fn

	outp.save(ofpptx);




def make_esp_trans(book, number):
	song, pathname, basename, rawname = get_song_paths(book, number)
	infile = "ehsf/esp/" + book + "/bil/" + book + "-" + song + "-bil.pptx";

	meta = dict()

	# extract metadata from first slide
	prs = Presentation(infile)
	meta_slide = prs.slides[0]
	for ndy, shape in enumerate(meta_slide.shapes, 1):
		if shape.shape_id == 2:
			meta['title'] = shape.text_frame.text
		elif shape.shape_id == 3:
			meta['credits'] = shape.text_frame.text
		elif shape.shape_id == 4:
			if len(shape.text_frame.text) != 0:
				meta['translation_copyright'] = shape.text_frame.text

	# esp/pftl/bil/xyz-bil.pptx (deck with embedded text boxes for subtitles)
	# Per song directories
	# esp/pftl/bil/xyz/*.png, *.txt (intermediate output from exporting from pptx file)
	# esp/pftl/xyz/ -- json and pngs

	bilpath = '/'.join(['ehsf', 'esp', book, 'bil', song])

	files = glob.glob(bilpath + "/*.png")
	files.sort()

	pprint.pprint(bilpath)
	pprint.pprint(files)

	# Analyze the extent of all the images
	crop = dict()
	for ndx, file in enumerate(files, 1):
		with open(file, 'rb') as f:
			crop[ndx] = analyze_image(f)

	window, padding = set_crop_window(crop, meta)

	esppath = '/'.join(['ehsf', 'esp', book, song])
	espbase = esppath + "/" + book + "-" + song

	for ndx, file in enumerate(files, 1):
		with open(file, 'rb') as f:
			size_image_to_window(f, crop[ndx], window, padding, espbase, None, ndx)
		
	# Output JSON data
	pprint.pprint(meta)
	jsonpath = '/'.join(['ehsf', 'esp', book, song])
	with open(jsonpath + "/" + book + "-" + song + ".json", 'w') as jsonfile:
		json.dump(meta, jsonfile, ensure_ascii=False, indent=4)

#
# Support for manipulating existing songs
#

def do_raw2png(picture, padding, basename, ndx):
	filename = basename + "-" + f"{ndx:02d}" + ".png"
	img = PIL.Image.open(picture)
	if img.mode != "RGB":
		img = img.convert(mode="RGB")

	# drw: need to set for each song
	mw = 2339
	mh = 1350
	box = (31.0, 219.0, 2370.0, 1569.0)
	print(box)
	img = img.crop(box)

	pw = mw * padding
	ph = mh * padding
	dx = (pw - mw) / 2
	dy = (ph - mh) / 2
	print(dx, dy, pw, ph)

	out = PIL.Image.new("RGB", (int(pw + 0.5), int(ph + 0.5)), color="white")
	out.paste(img, box=(int(dx), int(dy)))
	print(filename)
	out.save(filename)


def raw2png(book, number):
	song, pathname, basename, rawname = get_song_paths(book, number)

	rawpath = '/'.join(['ehsf', book, song, 'raw'])

	files = glob.glob(rawpath + "/*.png")
	files.sort()

	pprint.pprint(rawpath)
	pprint.pprint(files)

	for ndx, file in enumerate(files, 1):
		with open(file, 'rb') as f:
			do_raw2png(f, 1.05, basename, ndx)
		
#
# Support for custom (EH) songs
#

#
# Extracts images from PPT and exports to eh/123/raw/*.png
#
def raw2eh(number):
	song, pathname, basename, rawname = get_song_paths("eh", number)
	prs = Presentation("ehsf/eh/pptx/" + song + '-raw.pptx')

	# Create a directory to hold our output
	os.makedirs(pathname, exist_ok=True)
	os.makedirs(pathname + "/raw-input", exist_ok=True)

	slides = prs.slides

	for ndx, slide in enumerate(slides, 1):
		for ndy, shape in enumerate(slide.shapes, 1):
			print(ndy, shape.shape_type)
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#				filename = basename + "-" + f"{ndx:02d}" + ".png"
				img = PIL.Image.open(io.BytesIO(shape.image.blob))
				box = (0, 0, img.size[0], img.size[1] - 200)
# drw: crop out bottom				box = (0, 0, img.size[0], img.size[1] - 125)
#				box = (0, 0, img.size[0], img.size[1])
				img = img.crop(box)
				img.save(rawname + "-" + f"{ndx:02d}" + ".png")	# save the raw image

def process_eh_song(number):
	song, pathname, basename, rawname = get_song_paths("eh", number)

	# Create a directory to hold our output
	os.makedirs(pathname, exist_ok=True)
	os.makedirs(pathname + "/raw", exist_ok=True)

	# find all the raw files
	pngfiles = glob.glob(pathname + "/cooked/*.png")
	pngfiles.sort()
	pprint.pprint(pngfiles)

	# Analyze the extent of all the images
	meta = dict()
	crop = dict()
	for ndx, fn in enumerate(pngfiles):
		with open(fn, 'rb') as f:
			crop[ndx] = analyze_image(f)

	window, padding = set_crop_window(crop, meta)

	for ndx, fn in enumerate(pngfiles):
		with open(fn, 'rb') as f:
			size_image_to_window(f, crop[ndx], window, padding, basename, rawname, ndx+1)

	# Output JSON data
	pprint.pprint(meta)
	with open(basename + ".json", 'w') as jsonfile:
		json.dump(meta, jsonfile, ensure_ascii=False, indent=4)



def print_lyrics(book, number):
	song, paths = get_song_paths_new(book, number)
	with open(paths['engbase'] + ".json", 'r') as jsonfile:
		meta = json.load(jsonfile)
#	pprint.pprint(meta['lyrics'])
	print(book + "-" + song, meta['title'])
	print(meta['credits'])
	print(" ")
	for vs, text in meta['lyrics'].items():
		print("Verse", vs)
		for line in text:
			print(line)
		print(" ")


def before_brigham(espjson):
	modtime = os.path.getmtime(espjson)
	brigham = datetime.fromisoformat('2022-06-15T12:00:00')
#	print(modtime, datetime.fromtimestamp(modtime).isoformat())
	if modtime > datetime.timestamp(brigham):
		return False
	return True

def print_esp_list(brigham):
	# find all the raw files
	song_count = 0
	for book in ['pftl', 'phss', 'eh']:
		pathname = '/'.join(['ehsf', 'esp', book])
		files = glob.glob(pathname + "/*")
		files.sort()
		for file in files:
			xf = file.split('\\')	# drw: windows dependency
			number = xf[1]
			if not number.isnumeric():
				continue
#			if number == "bil" or number == "eng":
#				continue
			song, paths = get_song_paths_new(book, int(number))
			espjson = paths['espbase'] + ".json"
			if os.path.exists(espjson):
				needswork = before_brigham(espjson)
				good = ' ' if needswork else '*'
				if brigham and needswork:
					continue
				with open(paths['engbase'] + ".json", 'r') as jsonfile:
					meta = json.load(jsonfile)
					with open(espjson, 'r') as jsonesp:
						metaesp = json.load(jsonesp)
						print(good, book, song, meta['title'] + ' / ' + metaesp['title'])
						song_count = song_count + 1
	print("Number of translated songs =", song_count)



def export_json_metadata():
	# find all the raw files
	export = dict()
	for book in ['pftl', 'phss', 'eh']:
		pathname = '/'.join(['ehsf', book])
		files = glob.glob(pathname + "/[0-9][0-9][0-9]")	# only pure numbers
		files.sort()
		bookjson = dict()
		for file in files:
			xf = file.split('\\')	# drw: windows dependency
			number = xf[1]
			song, paths = get_song_paths_new(book, int(number))
			engjson = paths['engbase'] + ".json"
			espjson = paths['espbase'] + ".json"
			with open(engjson, 'r') as jsonfile:
				meta = json.load(jsonfile)
			if os.path.exists(espjson):
				with open(espjson, 'r') as jsonfile:
					esp = json.load(jsonfile)
					meta['esp'] = esp
			bookjson[number] = meta
		export[book] = bookjson

	with open("index.json", 'w') as outputfile:
		json.dump(export, outputfile, ensure_ascii=True, indent=4)


##############################################################################
##
## Song History Database
##
##############################################################################

def get_isodate(service, key):
	xdate = parser.parse(service[key])
	if 'time' in service:
		xtime = parser.parse(service['time'])
		h = xtime.hour
		m = xtime.minute
	else:
		if service['service'] == "Sunday Morning":
			h = 10
			m = 0
		elif service['service'] == "Sunday Evening":
			h = 18
			m = 0
		else:
			h = 19
			m = 30
	ndate = datetime(xdate.year, xdate.month, xdate.day, h, m, 0)
	return ndate.isoformat()

def parse_isodate(service):
	if 'isodate' in service:
		if len(service['isodate']) <= 10:
			date = get_isodate(service, 'isodate')
		else:
			date = service['isodate']
	elif 'date' in service:
		date = get_isodate(service, 'date')
	else:
		date = file		#"0000-00-00"
	return date


def add_leader_to_db(db, leader):
	lkey = None
	if leader != "unknown":
		# insert leader
		insert_leader = """INSERT OR IGNORE INTO leaders
		(first_name, last_name) 
		VALUES (?, ?);"""

		names = leader.split()
		leader_t = (names[0], names[1])

		cur = db.cursor()
		cur.execute(insert_leader, leader_t)
		db.commit()

		cur.execute("SELECT leader_id FROM leaders WHERE first_name = ? AND last_name = ?", leader_t)
		lkey = cur.fetchone()[0]

	return lkey




def add_service_type_to_db(db, service_type):
	insert_st = "INSERT OR IGNORE INTO service_types (description) VALUES (?);"

	cur = db.cursor()
	cur.execute(insert_st, (service_type, ))
	db.commit()

	cur.execute("SELECT service_type_id FROM service_types WHERE description = ?", (service_type,))
	key = cur.fetchone()[0]

	return key


def insert_service_to_db(db, service_type, date):
	stkey = add_service_type_to_db(db, service_type)

	insert_sv = "INSERT OR IGNORE INTO services (service_type_id, datetime) VALUES (?, ?);"

	cur = db.cursor()
	cur.execute(insert_sv, (stkey, date))
	db.commit()

	cur.execute("SELECT service_id FROM services WHERE datetime = ?", (date,))
	key = cur.fetchone()[0]

	return key


def get_service_type(service, date):
	service_map = { 'Sun - AM': 'Sunday Morning', 'Sun - EarlyAM': 'Sunday Morning', 'Sun - PM': 'Sunday Evening', 'Wed': 'Wednesday Evening', 'Wed - PM': 'Wednesday Evening', 'Gospel Meeting': 'Gospel Meeting', 'Sat - AM': 'Saturday Morning' }
	if 'service' in service:
		stype = service['service']	# Sunday Morning, Sunday Evening, Wednesday Evening, Gospel Meeting
		if stype.find('8:15') != -1:
			stype = 'Sunday Morning'
		elif stype.find('8:30') != -1:
			stype = 'Sunday Morning'
		elif stype.find('10:00') != -1:
			stype = 'Sunday Morning'
		elif stype.find('11:30') != -1:
			stype = 'Sunday Morning'
	elif 'type' in service:
		stypex = service['type']		# 'Sun - AM', 'Sun - PM', 'Wed', 'Sun - EarlyAM', 'Gospel Meeting'
		if stypex in service_map:
			stype = service_map[stypex]
		else:
			stype = stypex
	else:
		xdate = parser.parse(date)
		if xdate.hour < 12:
			stype = 'Sunday Morning'
		elif xdate.hour == 18:
			stype = 'Sunday Evening'
		elif xdate.hour == 19:
			stype = 'Wednesday Evening'
		else:
			stype = 'unknown'
	return stype


def get_verse_string(item):
	book = get_item(item, 'book')
	number = int(get_item(item, 'song'))

	song, paths = get_song_paths_new(book, number)
	with open(paths['engbase'] + ".json", 'r') as jsonfile:
		meta = json.load(jsonfile)
	custom = paths['engbase'] + "-custom.json"
	if os.path.exists(custom):
		with open(custom, 'r') as jsonfile:
			meta.update(json.load(jsonfile))

	verses = get_item(item, 'verses')
	chorus = get_item(item, 'chorus')
	coda = get_item(item, 'coda')
	if coda is None:
		coda = False

	wp = get_song_bubbles(meta, verses, chorus, coda)	# tbd: handle codas

	verse_string = ''
	for bubble in wp:
		if bubble["sing"]:
			verse_string += bubble["bubble"]

	return verse_string


insert_song = "INSERT INTO songs (book, songnum, verses, position, usage_id, service_id, leader_id) VALUES (?, ?, ?, ?, ?, ?, ?);"

def add_service_to_db(db, file):
	with open(file, 'r') as jsonfile:
		service = json.load(jsonfile)
		if 'skipdb' in service:
			return
		if 'leader' in service:
			sleader = service['leader']
		else:
			sleader = "unknown"
		date = parse_isodate(service)
		service_type = get_service_type(service, date)
		skey = insert_service_to_db(db, service_type, date)

		if 'items' in service:
			pos = 1
			for item in service['items']:
				if item['type'] == 'song' or item['type'] == 'song-music':
					if 'leader' in item:
						leader = item['leader']
					else:
						leader = sleader
					print(date, item['book'], item['song'], leader, pos)
					lkey = add_leader_to_db(db, leader)

					verse_string = get_verse_string(item)
					print(item['book'], int(item['song']), verse_string)

					song_data = (item['book'], int(item['song']), verse_string, pos, None, skey, lkey)

					cur = db.cursor()
					cur.execute(insert_song, song_data)
					db.commit()

					pos = pos + 1


def build_database():
	db = sqlite3.connect("test.db")

	years = glob.glob("worship/2[0-9][0-9][0-9]")	# only pure numbers
	years.sort()
	for year in years:
		files = glob.glob(year + "/*.json")
		files.sort()
		for file in files:
			add_service_to_db(db, file)
			print(' ')

def query_database(args):
	db = sqlite3.connect("test.db")
	cur = db.cursor()

	sql = "SELECT * FROM songs JOIN services USING(service_id) LEFT JOIN leaders USING (leader_id) WHERE book =? AND songnum =?"
	cur.execute(sql, (args.book, args.song))
	rows = cur.fetchall()
	for row in rows:
		print(row)

##############################################################################
##
## Main Processing
##
##############################################################################

def main():
	book_list = ['pftl', 'shs', 'phss', 'eh', "phss-eh"]
	lang_list = ['eng', 'esp', 'bil']
	parser = argparse.ArgumentParser(description="Song Slide Tool")
	parser.add_argument('-b', dest='book', help="Song book", choices=book_list, default=book_list[0])
	parser.add_argument('-s', dest='song', help="Song number", default="1")
	parser.add_argument('-w', dest='worship', help="Input worship JSON file", default=None)
	parser.add_argument('-d', dest='slide', help="esp-eng slide number to process", default=None)
	parser.add_argument('-lang', dest='language', help="Output language", choices=lang_list, default=lang_list[0])
	parser.add_argument('-output', dest='output', help="Output filename", default="output.pptx")
	parser.add_argument('--ppt', dest='ppt', help="Convert EHSF song to PowerPoint", action='store_true')
	parser.add_argument('--esp-eng', dest='espblank', help="Generate English deck for translation", action='store_true')
	parser.add_argument('--esp-bil', dest='esptrans', help="Process Bilingual deck into EHSF", action='store_true')
	parser.add_argument('--esp-list', dest='esplist', help="List available bilingual songs", action='store_true')
	parser.add_argument('--brigham', dest='brigham', help="Only since Brigham reviewed", action='store_true')
	parser.add_argument('--raw', dest='raw2png', help="Process raw files to png", action='store_true')
	parser.add_argument('--eh', dest='raw2eh', help="Process raw files to eh", action='store_true')
	parser.add_argument('--lyrics', dest='lyrics', help="Print lyrics of a song", action='store_true')
	parser.add_argument('--export', dest='export', help="Export JSON metadata", action='store_true')
	parser.add_argument('--db', dest='db', help="Build history database", action='store_true')
	parser.add_argument('--hist', dest='hist', help="Query history database", action='store_true')
	args = parser.parse_args()

	if args.worship is not None:
		make_worship_deck(args.worship)
	elif args.ppt:
		make_deck(args.book, int(args.song), args.language, args.output)
	elif args.espblank:
		make_esp_blank(args.book, int(args.song), args.slide)
	elif args.esptrans:
		make_esp_trans(args.book, int(args.song))
	elif args.esplist:
		print_esp_list(args.brigham)
	elif args.raw2png:
		raw2png(args.book, int(args.song))
	elif args.raw2eh:
		raw2eh(int(args.song))
	elif args.lyrics:
		print_lyrics(args.book, int(args.song))
	elif args.export:
		export_json_metadata()
	elif args.db:
		build_database()
	elif args.hist:
		query_database(args)
	else:
		if args.book == "pftl":
			process_pftl_song(int(args.song))
		elif args.book == "phss":
			process_phss_song_ppt(int(args.song))
		elif args.book == "eh":
			process_eh_song(int(args.song))
		elif args.book == "phss-eh":
			process_phss_to_eh(int(args.song))
		else:
			print("book not yet supported")


if __name__ == '__main__':
	main()

